"""학위논문 초안 요약 발표 슬라이드 (.pptx) 생성기.

대상: 학위논문 심사 / 박운상 지도교수 미팅
형식: 16:9, 25 슬라이드, 한국어, Malgun Gothic
산출: paper/thesis_slides.pptx

구조 (4부):
  1부 도입 (5장): 표지 / 목차 / 배경 / 문제 제기 / RQ + 기여
  2부 방법 (7장): 시스템 / 데이터 / 예측 / 해석 / Fusion / LLM-RAG / 평가+공정성
  3부 결과 (8장): 성능 / SHAP×Attn / 환각 / NLI / G-Eval / Persona / 공정성 / 일반화
  4부 결론 (5장): 핵심 발견 / 한계 / 응용 가이드 / 향후 연구 / Q&A

실행:
  PYTHONIOENCODING=utf-8 .venv/Scripts/python.exe -c \
    "import sys; sys.path.insert(0, 'thesis/draft'); import gen_thesis_slides; gen_thesis_slides.main()"
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Inches, Pt
from lxml import etree

# ─────────────────────────────────────────────────────────────
# 경로
# ─────────────────────────────────────────────────────────────
DRAFT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = DRAFT_DIR.parent.parent  # thesis/draft → thesis → project root
PAPER_DIR = PROJECT_ROOT / "paper"
FIG_DIR = PROJECT_ROOT / "figures"
PAPER_DIR.mkdir(parents=True, exist_ok=True)

# ─────────────────────────────────────────────────────────────
# 색상 팔레트 (gen_slides.py와 동일)
# ─────────────────────────────────────────────────────────────
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)
COLOR_ACCENT = RGBColor(0xC4, 0x4E, 0x52)
COLOR_HIGHLIGHT = RGBColor(0x55, 0xA8, 0x68)
COLOR_SUB = RGBColor(0x55, 0x55, 0x55)
COLOR_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KOREAN_FONT = "Malgun Gothic"
TOTAL = 27


# ─────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────
def _set_korean_font(run, size: int = 20, bold: bool = False,
                       color: RGBColor | None = None,
                       font_name: str = KOREAN_FONT):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ea_old in rPr.findall(qn("a:ea")):
        rPr.remove(ea_old)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    for la_old in rPr.findall(qn("a:latin")):
        rPr.remove(la_old)
    la = etree.SubElement(rPr, qn("a:latin"))
    la.set("typeface", font_name)


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, left, top, width, height,
                  text: str, size: int = 20, bold: bool = False,
                  color: RGBColor | None = None,
                  align: str = "left", anchor: str = "top"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    _set_korean_font(run, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, left, top, width, height, bullets: list,
                  size: int = 16, color: RGBColor | None = None):
    """bullets: 문자열 또는 dict {text, bold, color, sub:[]}."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if isinstance(b, dict):
            text = b["text"]
            sub = b.get("sub", [])
            b_bold = b.get("bold", False)
            b_color = b.get("color", color)
        else:
            text = b
            sub = []
            b_bold = False
            b_color = color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + text
        _set_korean_font(run, size=size, bold=b_bold, color=b_color)
        for s in sub:
            sp = tf.add_paragraph()
            sp.level = 1
            sp.alignment = PP_ALIGN.LEFT
            sr = sp.add_run()
            sr.text = "    –  " + s
            _set_korean_font(sr, size=size - 2, color=COLOR_SUB)
    return tb


def add_header_bar(slide, slide_w: int, title: str,
                     subtitle: str | None = None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(0), Cm(0), Cm(0.5), Cm(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_textbox(slide, Cm(1.0), Cm(0.4), Cm(30), Cm(1.0),
                  title, size=24, bold=True, color=COLOR_PRIMARY)
    if subtitle:
        add_textbox(slide, Cm(1.0), Cm(1.4), Cm(30), Cm(0.8),
                      subtitle, size=12, color=COLOR_SUB)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(1.0), Cm(2.4),
                                       Cm(31.5), Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()


def add_footer(slide, slide_w: int, slide_h: int, page_no: int):
    add_textbox(slide, Cm(1.0), Cm(slide_h / 360000 - 0.9),
                  Cm(22), Cm(0.6),
                  "TabNet × SHAP × LLM XAI-RAG  |  A70067 오현택",
                  size=9, color=COLOR_SUB)
    add_textbox(slide, Cm(slide_w / 360000 - 3.0),
                  Cm(slide_h / 360000 - 0.9),
                  Cm(2), Cm(0.6), f"{page_no} / {TOTAL}",
                  size=9, color=COLOR_SUB, align="right")


def add_picture_centered(slide, path: Path, top_cm: float,
                            height_cm: float, slide_width: int):
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), Cm(0), Cm(top_cm),
                                       height=Cm(height_cm))
    pic.left = (slide_width - pic.width) // 2
    return pic


def add_table(slide, left, top, col_widths, row_heights, rows,
                header: bool = True, header_color: RGBColor = COLOR_PRIMARY,
                font_size: int = 12):
    """간단 표 — 도형 기반."""
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            x = left + sum(col_widths[:ci], Emu(0))
            y = top + sum(row_heights[:ri], Emu(0))
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                                col_widths[ci],
                                                row_heights[ri])
            if header and ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = header_color
                fg = COLOR_WHITE
                fb = True
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = (RGBColor(0xF8, 0xF9, 0xFB)
                                                  if ri % 2 == 1 else
                                                  COLOR_WHITE)
                fg = COLOR_DARK
                fb = False
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.margin_right = Cm(0.1)
            tf.margin_top = Cm(0.08)
            tf.margin_bottom = Cm(0.08)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell)
            _set_korean_font(run, size=font_size, bold=fb, color=fg)


def add_callout_box(slide, left, top, width, height, text: str,
                      size: int = 14, color: RGBColor = COLOR_HIGHLIGHT,
                      bold: bool = True):
    """핵심 메시지 강조용 색상 박스."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_korean_font(run, size=size, bold=bold, color=COLOR_WHITE)


# ─────────────────────────────────────────────────────────────
# 슬라이드 정의
# ─────────────────────────────────────────────────────────────
def make_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    SH_CM = SH / 360000

    # ═══════════════════════════════════════════════════
    # Slide 1 — 표지
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    bg.line.fill.background()

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(0), Cm(0), Cm(1.5), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_textbox(s, Cm(2.5), Cm(3.5), Cm(28), Cm(3.5),
                  "정형 데이터 특화 딥러닝(TabNet)과\n"
                  "거대언어모델(LLM) 기반 XAI-RAG를 활용한\n"
                  "설명 가능한 신용 평가 및 사용자 맞춤형 리포트 생성",
                  size=26, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Cm(2.5), Cm(8.0), Cm(28), Cm(1.0),
                  "— 어텐션-SHAP 융합 컨텍스트와 다중 평가 프레임워크 적용 —",
                  size=16, color=COLOR_SUB)
    add_textbox(s, Cm(2.5), Cm(11.0), Cm(28), Cm(1.0),
                  "석사 학위논문 초안 발표",
                  size=18, bold=True, color=COLOR_ACCENT)
    add_textbox(s, Cm(2.5), Cm(13.5), Cm(28), Cm(3.0),
                  "데이터사이언스 · 인공지능 전공\n"
                  "학번 A70067  |  오현택\n"
                  "지도교수: 박운상\n"
                  "서강대학교 AI · SW대학원  |  2026년 6월",
                  size=14, color=COLOR_DARK)

    # ═══════════════════════════════════════════════════
    # Slide 2 — 목차
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "목 차", "Contents")

    parts = [
        ("1부. 도입", "연구 배경 · 문제 제기 · 연구 질문 · 기여",
            COLOR_PRIMARY),
        ("2부. 연구 방법", "시스템 구조 · 데이터 · 예측·해석 · Fusion · LLM-RAG · 평가",
            COLOR_HIGHLIGHT),
        ("3부. 분석 결과", "예측 성능 · 해석 일관성 · 4-mode 비교 · Persona · 공정성 · 일반화",
            COLOR_ACCENT),
        ("4부. 결론 및 시사점", "핵심 발견 · 한계 · 응용 가이드 · 향후 연구",
            RGBColor(0x8C, 0x6B, 0xB1)),
    ]
    for i, (title, sub, color) in enumerate(parts):
        y_top = Cm(3.5 + i * 3.3)
        # 번호 배지
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Cm(2.0), y_top, Cm(1.8), Cm(1.8))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        brun = bp.add_run()
        brun.text = str(i + 1)
        _set_korean_font(brun, size=22, bold=True, color=COLOR_WHITE)
        # 제목 + 부제
        add_textbox(s, Cm(4.5), y_top, Cm(25), Cm(1.0),
                      title, size=20, bold=True, color=COLOR_PRIMARY)
        add_textbox(s, Cm(4.5), y_top + Cm(1.0), Cm(25), Cm(1.0),
                      sub, size=13, color=COLOR_SUB)
    add_footer(s, SW, SH, 2)

    # ═══════════════════════════════════════════════════
    # Slide 3 — 연구 배경
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "1. 연구 배경",
                       "신용 평가에 \"설명 가능한 자연어\"가 필요한 이유")
    add_bullets(s, Cm(1.0), Cm(3.2), Cm(31), Cm(14), [
        {"text": "신용 평가 자동화와 블랙박스 문제", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
            "XGBoost·LightGBM·TabNet — 부도 예측 성능을 크게 끌어올림",
            "그러나 내부 의사결정이 직관적이지 않은 \"블랙박스\" 특성"]},
        {"text": "제도적·윤리적 설명 책임", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
            "EU GDPR 제22조 — 자동화된 의사결정에 대한 설명 받을 권리",
            "금융위·금감원 「금융 분야 AI 가이드라인」 (2021) — 명시적 설명 의무",
            "성별·연령 기반 잠재적 편향 → 차별 가능성"]},
        {"text": "기존 XAI(SHAP·LIME)의 한계", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
            "변수 중요도 그래프·기여도 수치 → 전문가 외엔 해석 어려움",
            "\"모델은 설명 가능해졌으나 사람에게는 여전히 설명되지 않는다\""]},
        {"text": "LLM-RAG의 가능성과 위험", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
            "데이터-투-텍스트 변환에 강함, RAG로 사실성 보강 가능",
            "그러나 환각(Hallucination) 문제 — 신용평가는 정확성과 법적 책임이 중요"]},
    ], size=14)
    add_footer(s, SW, SH, 3)

    # ═══════════════════════════════════════════════════
    # Slide 4 — 문제 제기 (5가지 한계)
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "2. 문제 제기",
                       "기존 연구의 5가지 학술적 공백")

    limitations = [
        ("①", "TabNet 어텐션 활용 부재",
            "TabNet은 주로 단독 예측 모델로만 사용 — 어텐션 마스크와 SHAP를 결합한 사례 거의 없음"),
        ("②", "LLM 환각 차단의 차별성 검증 부재",
            "단순 hard constraint만으로 환각 차단 가능 — SHAP 기반 RAG의 추가 가치는?"),
        ("③", "공정성 보정의 신용평가 통합 부재",
            "Reweighing 등 정식 mitigation을 모델 성능 손실 최소화로 통합 적용한 사례 부족"),
        ("④", "단일 데이터셋의 일반화 한계",
            "선행 연구 대부분 단일 데이터셋에서만 평가 — 데이터셋 의존성 / 일반 패턴 구별 어려움"),
        ("⑤", "이해관계자별 평가 부재",
            "신용 전문가·고객·규제기관 등 페르소나 관점의 다중 평가 프레임워크 미정립"),
    ]
    for i, (num, title, desc) in enumerate(limitations):
        y = Cm(3.4 + i * 2.6)
        # 번호 배지
        add_textbox(s, Cm(1.0), y, Cm(1.5), Cm(2.0),
                      num, size=28, bold=True, color=COLOR_ACCENT)
        # 제목 + 설명
        add_textbox(s, Cm(2.6), y, Cm(29), Cm(1.0),
                      title, size=16, bold=True, color=COLOR_PRIMARY)
        add_textbox(s, Cm(2.6), y + Cm(0.9), Cm(29), Cm(1.5),
                      desc, size=12, color=COLOR_DARK)
    add_footer(s, SW, SH, 4)

    # ═══════════════════════════════════════════════════
    # Slide 5 — 연구 질문 + 기여
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "3. 연구 질문 (RQ) 및 기여 (C)",
                       "4 개의 연구 질문, 5 개의 학술적 기여")

    # 왼쪽 — RQ
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(15), Cm(0.8),
                  "연구 질문 (Research Questions)", size=16, bold=True,
                  color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.0), Cm(4.0), Cm(15.5), Cm(13), [
        {"text": "RQ1. SHAP × TabNet 어텐션 동의 기반 융합 컨텍스트가 "
                  "LLM 자연어 설명의 사실 기반 충실성을 향상시키는가?",
            "color": COLOR_DARK},
        {"text": "RQ2. Hard constraint만으로 환각 차단이 충분한 상황에서, "
                  "Fusion 컨텍스트의 차별적 가치는?",
            "color": COLOR_DARK},
        {"text": "RQ3. Reweighing 기반 사전처리 공정성 보정이 "
                  "성능 손실 없이 4/5 규칙을 통과시키는가?",
            "color": COLOR_DARK},
        {"text": "RQ4. Fusion 메커니즘이 다른 데이터셋(UCI German)에도 "
                  "일관 작동하며 Persona별 trade-off는?",
            "color": COLOR_DARK},
    ], size=12)

    # 오른쪽 — C
    add_textbox(s, Cm(17.0), Cm(3.2), Cm(15), Cm(0.8),
                  "본 연구의 기여 (Contributions)", size=16, bold=True,
                  color=COLOR_HIGHLIGHT)
    add_bullets(s, Cm(17.0), Cm(4.0), Cm(15.5), Cm(13), [
        {"text": "C1. Agreement-aware Fusion Context 신규 설계 "
                  "(agreed · shap_only · attention_only)",
            "color": COLOR_DARK},
        {"text": "C2. 4-tier 다중 평가 프레임워크 구축 "
                  "(Rules + NLI + Cross-Judge G-Eval + Persona Pilot)",
            "color": COLOR_DARK},
        {"text": "C3. 4-mode 비교를 통한 차별성 입증 — "
                  "환각 차단 ≠ fact-grounded 정확성",
            "color": COLOR_DARK},
        {"text": "C4. 두 데이터셋(Home Credit, UCI German) 일반화 검증",
            "color": COLOR_DARK},
        {"text": "C5. Reweighing 공정성 보정 — 4/4 조합 모두 4/5 규칙 통과",
            "color": COLOR_DARK},
    ], size=12)
    add_footer(s, SW, SH, 5)

    # ═══════════════════════════════════════════════════
    # Slide 6 — 시스템 전체 구조
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "4. 시스템 전체 구조",
                       "4 단계 파이프라인 — 예측 → 해석 → 융합 → 자연어 + 평가")

    # 파이프라인 그림 (없으면 박스로 대체)
    pipe_path = FIG_DIR / "42_thesis_pipeline.png"
    if pipe_path.exists():
        add_picture_centered(s, pipe_path, top_cm=3.2, height_cm=10.0,
                                slide_width=SW)
    else:
        # 4 단계 박스 fallback
        stages = [
            ("① 예측", "XGBoost / TabNet\nP(default) 산출",
                RGBColor(0x4C, 0x72, 0xB0)),
            ("② 해석", "SHAP local +\nTabNet 어텐션",
                RGBColor(0xDD, 0x88, 0x52)),
            ("③ 융합 컨텍스트", "Agreement-aware\nJSON + 마스킹",
                COLOR_HIGHLIGHT),
            ("④ LLM-RAG + 평가", "자연어 설명 +\n4-tier 평가",
                RGBColor(0x8C, 0x6B, 0xB1)),
        ]
        box_w = Cm(7); box_h = Cm(4.5); gap = Cm(1.0)
        total_w = 4 * box_w + 3 * gap
        start_x = (SW - total_w) // 2
        box_y = Cm(5.5)
        for i, (label, desc, color) in enumerate(stages):
            x = start_x + i * (box_w + gap)
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          x, box_y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            tf = box.text_frame
            tf.margin_left = Cm(0.3); tf.margin_top = Cm(0.3)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run(); r1.text = label
            _set_korean_font(r1, size=18, bold=True, color=COLOR_WHITE)
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = desc
            _set_korean_font(r2, size=12, color=COLOR_WHITE)

    add_callout_box(s, Cm(2.0), Cm(15.5), Cm(29.5), Cm(1.5),
                      "★ 본 연구 핵심: ③ Agreement-aware Fusion Context — "
                      "두 해석 정보의 동의 여부 자체를 LLM 프롬프트 신호로 활용",
                      size=14, color=COLOR_HIGHLIGHT)
    add_footer(s, SW, SH, 6)

    # ═══════════════════════════════════════════════════
    # Slide 7 — 데이터셋
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "5. 데이터셋",
                       "Home Credit (메인) + UCI German Credit (일반화 검증)")

    rows = [
        ["항목", "Home Credit Default Risk", "UCI German Credit"],
        ["출처", "Kaggle Competition (2018)", "UCI ML Repository (1994)"],
        ["관측치", "307,511 명", "1,000 명"],
        ["원본 변수", "122 개", "21 개"],
        ["전처리 후 변수", "214 개 (one-hot + flag)", "63 개 (one-hot + flag)"],
        ["타깃", "TARGET (부도 여부)", "credit-risk (good / bad)"],
        ["부도율", "약 8.07 %", "30.0 %"],
        ["분할", "60 / 20 / 20 stratified (SEED=42)",
            "60 / 20 / 20 stratified (SEED=42)"],
        ["보호 속성", "CODE_GENDER, DAYS_BIRTH(연령)",
            "personal_status(성), age, foreign_worker"],
    ]
    add_table(s, Cm(1.5), Cm(3.4),
                [Cm(6.5), Cm(12.5), Cm(11.5)],
                [Cm(1.0)] * len(rows), rows, font_size=12)

    add_callout_box(s, Cm(2.0), Cm(15.5), Cm(29.5), Cm(1.5),
                      "두 데이터셋 — 동일 파이프라인 이식으로 메커니즘 일반화 검증",
                      size=14, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 7)

    # ═══════════════════════════════════════════════════
    # Slide 8 — 예측 모델
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "6. 예측 모델",
                       "XGBoost (메인) + TabNet (보조 + 어텐션 추출)")

    add_textbox(s, Cm(1.0), Cm(3.2), Cm(31), Cm(0.8),
                  "5-fold Stratified CV 성능 (Test set, mean ± std)",
                  size=15, bold=True, color=COLOR_PRIMARY)

    cv_rows = [
        ["모델", "AUROC (Home)", "AUROC (German)", "비고"],
        ["Logistic Regression", "0.7475 ± 0.0011", "0.7969 ± 0.0289",
            "전통 baseline"],
        ["LightGBM", "0.7574 ± 0.0009", "0.7656 ± 0.0312",
            "그래디언트 부스팅"],
        ["XGBoost (메인)", "0.7587 ± 0.0008", "0.7714 ± 0.0298",
            "본 연구 예측기"],
        ["TabNet", "0.7543 ± 0.0019", "0.7501 ± 0.0367",
            "어텐션 추출용"],
    ]
    add_table(s, Cm(1.5), Cm(4.5),
                [Cm(7.5), Cm(8.0), Cm(8.0), Cm(7.0)],
                [Cm(1.0)] * len(cv_rows), cv_rows, font_size=12)

    add_textbox(s, Cm(1.0), Cm(11.0), Cm(31), Cm(0.8),
                  "보조 테이블 추가 효과 (Home Credit)",
                  size=15, bold=True, color=COLOR_PRIMARY)
    aux_rows = [
        ["지표", "Baseline", "+ Bureau + PrevApp", "개선"],
        ["AUROC", "0.7587", "0.7755", "+2.22 %"],
        ["AUPRC", "0.2401", "0.2598", "+8.21 %"],
        ["KS 통계량", "0.4054", "0.4371", "+7.81 %"],
    ]
    add_table(s, Cm(1.5), Cm(12.2),
                [Cm(7.5), Cm(7.5), Cm(8.5), Cm(7.0)],
                [Cm(1.0)] * len(aux_rows), aux_rows, font_size=12)
    add_footer(s, SW, SH, 8)

    # ═══════════════════════════════════════════════════
    # Slide 9 — 모델 해석
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "7. 모델 해석",
                       "SHAP local 기여도 + TabNet 어텐션 마스크")

    # 왼쪽 — SHAP
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(15.5), Cm(1.0),
                  "SHAP (TreeSHAP)", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.0), Cm(4.2), Cm(15.5), Cm(10), [
        "변수의 부도 확률 기여도 (방향 + 크기)",
        "전체 변수에 대한 dense 기여도 — 풍부한 정보",
        "XGBoost 3.x base_score 파싱 버그 — _XgbNativeExplainer 우회 적용",
        "Top-K 부호 + 크기 정보를 LLM 프롬프트에 활용",
    ], size=13)

    # 오른쪽 — TabNet attention
    add_textbox(s, Cm(17.0), Cm(3.2), Cm(15.5), Cm(1.0),
                  "TabNet Attention", size=18, bold=True,
                  color=COLOR_ACCENT)
    add_bullets(s, Cm(17.0), Cm(4.2), Cm(15.5), Cm(10), [
        "Sparsemax 기반 변수 선택 — sparse 마스크",
        "인스턴스마다 다른 변수를 선택 (decision step 별)",
        "방향 정보 없음 — \"무엇을 봤는가\"만 표현",
        "SHAP과 다른 관점의 해석 정보 제공",
    ], size=13)

    add_callout_box(s, Cm(2.0), Cm(15.0), Cm(29.5), Cm(2.0),
                      "두 정보는 본질적으로 다른 해석 관점\n"
                      "(전체 dense 기여 vs sparse 변수 선택) — 융합이 정보 보완을 제공",
                      size=14, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 9)

    # ═══════════════════════════════════════════════════
    # Slide 10 — ★ Fusion Context
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "8. ★ Agreement-aware Fusion Context",
                       "본 연구의 핵심 — 동의 여부 자체를 LLM 신호로 활용")

    # 3 그룹 분해 박스
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(31), Cm(0.8),
                  "SHAP Top-K ∩ TabNet Attention Top-K → 3 그룹 분해",
                  size=15, bold=True, color=COLOR_PRIMARY)

    groups = [
        ("agreed", "양 모델이 모두 선택한 변수\n(두 해석이 동의)",
            "→ 가장 신뢰할 수 있는 신호로 LLM에 전달",
            COLOR_HIGHLIGHT),
        ("shap_only", "SHAP만 선택한 변수\n(부호+크기 정보 있음)",
            "→ 강한 dense 기여도 — 보조 신호",
            RGBColor(0x4C, 0x72, 0xB0)),
        ("attention_only", "TabNet attention만 선택한 변수\n(sparse 선택)",
            "→ \"모델이 본 변수\" — 추가 컨텍스트",
            RGBColor(0xDD, 0x88, 0x52)),
    ]
    box_w = Cm(10); box_h = Cm(7); gap = Cm(0.5)
    total_w = 3 * box_w + 2 * gap
    start_x = (SW - total_w) // 2
    box_y = Cm(4.3)
    for i, (name, desc, role, color) in enumerate(groups):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Cm(0.4); tf.margin_top = Cm(0.4)
        tf.margin_right = Cm(0.4)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = name
        _set_korean_font(r1, size=18, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + desc
        _set_korean_font(r2, size=12, color=COLOR_WHITE)
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = "\n" + role
        _set_korean_font(r3, size=11, color=COLOR_WHITE)

    add_callout_box(s, Cm(2.0), Cm(12.5), Cm(29.5), Cm(2.0),
                      "Agreement 통계 (n=100, Home): mean(agreed)=2.12, "
                      "n_agreed 분포 0~3 — 4+ 동의 0 %\n"
                      "보호 속성(GENDER, DAYS_BIRTH 등) JSON 단계에서 마스킹",
                      size=13, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 10)

    # ═══════════════════════════════════════════════════
    # Slide 11 — LLM-RAG 4-mode + Hard Constraints
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "9. LLM-RAG 4-Mode 비교 설계",
                       "Fusion의 차별성을 정량 입증하기 위한 ablation")

    modes = [
        ("no_shap", "원시 변수만\n(SHAP/RAG 없음)", "Baseline",
            RGBColor(0x99, 0x99, 0x99)),
        ("generic_rag", "원시 변수 + 7개 도메인 chunks\n(SHAP 없음)",
            "RAG ablation", RGBColor(0x4C, 0x72, 0xB0)),
        ("shaponly", "SHAP top-5만\n(attention 없음)",
            "SHAP-only", RGBColor(0xDD, 0x88, 0x52)),
        ("fusion", "SHAP + Attention\n동의 기반 융합 ★",
            "★ 본 연구", COLOR_HIGHLIGHT),
    ]
    box_w = Cm(7.3); box_h = Cm(6); gap = Cm(0.6)
    total_w = 4 * box_w + 3 * gap
    start_x = (SW - total_w) // 2
    box_y = Cm(3.8)
    for i, (name, desc, role, color) in enumerate(modes):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Cm(0.3); tf.margin_top = Cm(0.3)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = name
        _set_korean_font(r1, size=16, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + desc
        _set_korean_font(r2, size=11, color=COLOR_WHITE)
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = "\n" + role
        _set_korean_font(r3, size=11, bold=True, color=COLOR_WHITE)

    add_textbox(s, Cm(1.0), Cm(10.5), Cm(31), Cm(0.8),
                  "공통 Hard Constraints 5 항목 (모든 mode 적용)",
                  size=15, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.0), Cm(11.3), Cm(31), Cm(5), [
        "컨텍스트에 없는 변수명 생성 금지 (환각 차단)",
        "민감 속성(성별·연령) 언급 금지 — 마스킹된 컨텍스트만 사용",
        "값(value) 정확 인용 강제 — 임의 추정 금지",
        "5-section 출력 형식 (위험도·핵심 근거·세부 분석·추천·요약)",
        "LLM 비교: Anthropic Claude Sonnet 4.5 + Google Gemini 2.5 Flash",
    ], size=12)
    add_footer(s, SW, SH, 11)

    # ═══════════════════════════════════════════════════
    # Slide 12 — 4-tier 평가 + 공정성
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "10. 4-tier 평가 + 공정성 보정",
                       "Rules + NLI + G-Eval + Persona + Reweighing")

    tiers = [
        ("Tier 1", "룰 기반 평가", "변수·값·부호 토큰 매칭\nHallucination Rate (strict/broad)",
            RGBColor(0x4C, 0x72, 0xB0)),
        ("Tier 2", "NLI 기반 함의도",
            "mDeBERTa-v3 다국어 NLI\nEntailment / Contradiction 측정",
            RGBColor(0xDD, 0x88, 0x52)),
        ("Tier 3", "Cross-Judge G-Eval",
            "Claude × Gemini 양방향 평가자\n4 차원 (factual/completeness/sensitive/style)",
            COLOR_HIGHLIGHT),
        ("Tier 4", "Persona Pilot",
            "3 페르소나 (Expert/Customer/Regulator)\n3 metrics × 5 점 척도",
            RGBColor(0x8C, 0x6B, 0xB1)),
    ]
    box_w = Cm(7.3); box_h = Cm(5.5); gap = Cm(0.6)
    total_w = 4 * box_w + 3 * gap
    start_x = (SW - total_w) // 2
    box_y = Cm(3.8)
    for i, (tier, name, desc, color) in enumerate(tiers):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Cm(0.3); tf.margin_top = Cm(0.3)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = tier
        _set_korean_font(r0, size=11, color=COLOR_WHITE)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = name
        _set_korean_font(r1, size=16, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + desc
        _set_korean_font(r2, size=11, color=COLOR_WHITE)

    add_textbox(s, Cm(1.0), Cm(10.5), Cm(31), Cm(0.8),
                  "공정성 보정 — Kamiran-Calders Reweighing (pre-processing)",
                  size=15, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.0), Cm(11.3), Cm(31), Cm(5), [
        "보호 속성: CODE_GENDER, DAYS_BIRTH(AGE), personal_status",
        "정량 지표: Demographic Parity, Equal Opportunity, Disparate Impact (4/5 rule)",
        "비교 대상: Fairlearn ExpGrad (DP / EO) — DP는 AUROC 큰 손실, EO는 안티 패턴",
    ], size=13)
    add_footer(s, SW, SH, 12)

    # ═══════════════════════════════════════════════════
    # Slide 13 — 모델 성능 + AUROC 실무 평가
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "11. 예측 성능 — 신용평가 산업 벤치마크 대비",
                       "AUROC 0.76 = 실무 수준 (Lessmann et al. 2015 기준)")

    add_textbox(s, Cm(1.0), Cm(3.2), Cm(31), Cm(0.8),
                  "신용평가 도메인 AUROC 분류 기준 (Lessmann et al., 2015)",
                  size=14, bold=True, color=COLOR_PRIMARY)

    bench_rows = [
        ["AUROC 구간", "도메인 평가", "실무 의미"],
        ["≥ 0.85", "최상위 — 거의 보지 못함",
            "데이터 누수 또는 비현실적 가정 의심"],
        ["0.80 ~ 0.85", "상위 — 가장 좋은 모델", "특수 도메인 / 풍부한 데이터"],
        ["0.70 ~ 0.80", "Mainstream — 실무 표준",
            "★ 본 연구 (0.7587, 0.7714)"],
        ["0.60 ~ 0.70", "Borderline", "도입 가능하지만 신중 검토"],
        ["< 0.60", "Below standard", "재학습 필요"],
    ]
    add_table(s, Cm(1.5), Cm(4.2),
                [Cm(7), Cm(11), Cm(13)],
                [Cm(1.0)] * len(bench_rows), bench_rows, font_size=12)

    add_callout_box(s, Cm(2.0), Cm(11.5), Cm(29.5), Cm(2.5),
                      "본 연구 메시지: 예측 정확도 갱신이 아니라 "
                      "통합 시스템 구축\n"
                      "Baseline AUROC 0.76 (mainstream)을 honest 보고 → "
                      "Fusion 차별성은 NLI Entailment 0.625 + Value Match 0.88에서 정량 입증",
                      size=14, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 13)

    # ═══════════════════════════════════════════════════
    # Slide 14 — SHAP × Attention 일관성
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "12. SHAP × Attention 일관성",
                       "두 해석 정보의 부분 일관 + 부분 상보 — Fusion 정당성")

    sa_path = FIG_DIR / "16_attention_vs_shap_scatter.png"
    if sa_path.exists():
        add_picture_centered(s, sa_path, top_cm=3.5, height_cm=9.5,
                                slide_width=SW)

    add_callout_box(s, Cm(1.0), Cm(13.5), Cm(15.5), Cm(3.5),
                      "Home Credit (n=214)\n"
                      "Spearman ρ (full) = 0.117\n"
                      "Top-50 overlap ρ = -0.195",
                      size=14, color=COLOR_PRIMARY)
    add_callout_box(s, Cm(17.0), Cm(13.5), Cm(15.5), Cm(3.5),
                      "UCI German (n=63)\n"
                      "Spearman ρ (full) = 0.114\n"
                      "→ 두 데이터셋 거의 동일 일관성 ★",
                      size=14, color=COLOR_HIGHLIGHT)
    add_footer(s, SW, SH, 14)

    # ═══════════════════════════════════════════════════
    # Slide 15 — 환각률 0% 결과 + Hard Constraints 5항목
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "13. 환각률 0 % 결과 + Hard Constraints",
                       "모든 4-mode × 2 LLM × 데이터셋에서 환각 0건 — 그 메커니즘")

    # 좌측 — 환각률 결과 표
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(15.5), Cm(0.8),
                  "(1) 환각률 측정 결과 (300 호출, 양 데이터셋)",
                  size=13, bold=True, color=COLOR_PRIMARY)
    halluc_rows = [
        ["Mode", "Home", "German"],
        ["no_shap", "0 / 60", "0 / 60"],
        ["generic_rag", "0 / 60", "0 / 60"],
        ["shaponly", "0 / 60", "0 / 60"],
        ["fusion ★", "0 / 60", "0 / 60"],
    ]
    add_table(s, Cm(1.0), Cm(4.2),
                [Cm(6.5), Cm(4.5), Cm(4.5)],
                [Cm(0.95)] * len(halluc_rows), halluc_rows, font_size=12)

    add_callout_box(s, Cm(1.0), Cm(10.0), Cm(15.5), Cm(2.5),
                      "★ 모든 mode에서 0 %\n"
                      "→ Fusion의 차별성 근거가 아님\n"
                      "  단순 hard constraint만으로 충분",
                      size=12, color=COLOR_ACCENT)

    # 우측 — Hard Constraints 5항목 (메커니즘)
    add_textbox(s, Cm(17.0), Cm(3.2), Cm(15.5), Cm(0.8),
                  "(2) 0 % 달성 메커니즘 — Hard Constraints 5항목",
                  size=13, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(17.0), Cm(4.0), Cm(15.5), Cm(11), [
        {"text": "컨텍스트에 없는 변수명 생성 금지",
            "bold": True, "color": COLOR_DARK,
            "sub": ["LLM 출력의 영문 토큰 = 입력 변수 집합 ⊆"]},
        {"text": "민감 속성 언급 금지",
            "bold": True, "color": COLOR_DARK,
            "sub": ["성별·연령 등 마스킹된 컨텍스트만 사용"]},
        {"text": "값(value) 정확 인용 강제",
            "bold": True, "color": COLOR_DARK,
            "sub": ["임의 추정·반올림·생략 금지"]},
        {"text": "5-section 출력 형식 고정",
            "bold": True, "color": COLOR_DARK,
            "sub": ["위험도·핵심 근거·세부 분석·추천·요약"]},
        {"text": "컨텍스트 외 외부 데이터 인용 금지",
            "bold": True, "color": COLOR_DARK,
            "sub": ["LLM 내부 지식만으로 변수·값 생성 금지"]},
    ], size=11)

    add_callout_box(s, Cm(1.0), Cm(15.5), Cm(31.5), Cm(1.5),
                      "→ 다음 슬라이드: 환각률 0 %의 정확한 의미 — "
                      "3 가지 정확도 차원 분리 + Counterfactual 대비",
                      size=12, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 15)

    # ═══════════════════════════════════════════════════
    # Slide 16 — 환각률 0 %의 정확한 의미 (3차원 분리 + Counterfactual)
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "14. 환각률 0 %의 정확한 의미",
                       "세 가지 정확도 차원 분리 + Counterfactual 45.5 % → 0 % 대비")

    # 3 차원 분리 표
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(31), Cm(0.8),
                  "본 시스템의 3 가지 정확도 차원 — 환각률 0 %는 차원 B만 보장",
                  size=14, bold=True, color=COLOR_PRIMARY)

    dim_rows = [
        ["차원", "측정 대상", "측정 도구", "본 연구 결과"],
        ["A. 부도 예측 정확도",
            "XGBoost가 실제 부도/정상 맞춤",
            "AUROC · AUPRC · KS",
            "0.7587 (Home) / 0.7714 (German)"],
        ["B. LLM 변수명 환각률 ★",
            "LLM이 없는 변수명 생성 비율",
            "룰 기반 영문 토큰 매칭",
            "★ 0 % (모든 mode)"],
        ["C. LLM 의미적 충실성",
            "설명과 컨텍스트의 의미 일치",
            "NLI Entailment (mDeBERTa)",
            "fusion 0.625 / no_shap 0.30"],
    ]
    add_table(s, Cm(1.0), Cm(4.2),
                [Cm(7), Cm(9), Cm(7), Cm(8.5)],
                [Cm(1.4), Cm(1.4), Cm(1.4), Cm(1.4)],
                dim_rows, font_size=11)

    # 보장 / 미보장
    add_callout_box(s, Cm(1.0), Cm(10.8), Cm(15.5), Cm(3.5),
                      "✓ 0 %가 보장하는 것\n"
                      "LLM이 출력에 등장시킨 변수 토큰이 모두 입력 컨텍스트 또는 "
                      "데이터셋 변수 집합에 존재한다는 것",
                      size=11, color=COLOR_HIGHLIGHT)
    add_callout_box(s, Cm(17.0), Cm(10.8), Cm(15.5), Cm(3.5),
                      "✗ 0 %가 보장하지 않는 것\n"
                      "① 부도 예측이 실제와 일치 (차원 A) "
                      "② 값을 100 % 정확히 인용 (Value Match) "
                      "③ 설명 의미가 입력과 100 % 일치 (차원 C)",
                      size=11, color=COLOR_ACCENT)

    # Counterfactual 대비
    add_callout_box(s, Cm(1.0), Cm(15.0), Cm(31.5), Cm(2.0),
                      "★ Counterfactual 검증 (§4.3.6) — Hard Constraint 제거 시 "
                      "Gemini 환각률 45.5 % → 적용 후 0 %\n"
                      "즉 0 %는 무에서 나온 결과가 아니라 45.5 %를 차단한 정량 효과 "
                      "— Hard Constraint 정책의 견고함 입증",
                      size=12, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 16)

    # ═══════════════════════════════════════════════════
    # Slide 17 — ★ NLI Entailment 4-mode 단조증가
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "15. ★ NLI Entailment 4-mode 단조증가",
                       "Fusion 일관 1위 — 두 데이터셋에서 동일 패턴")

    add_textbox(s, Cm(1.0), Cm(3.3), Cm(31), Cm(0.8),
                  "NLI Entailment Rate (mDeBERTa-v3 multilingual)",
                  size=15, bold=True, color=COLOR_PRIMARY)

    nli_rows = [
        ["Mode", "Home Credit", "UCI German Credit",
            "Value Match Rate (Home)"],
        ["no_shap", "0.350", "0.393", "0.59"],
        ["generic_rag", "0.367", "0.410", "0.73"],
        ["shaponly", "0.461", "0.628", "0.85"],
        ["fusion ★", "0.625", "0.711", "0.90"],
    ]
    add_table(s, Cm(1.5), Cm(4.3),
                [Cm(7), Cm(9), Cm(9), Cm(6.5)],
                [Cm(1.1)] * len(nli_rows), nli_rows, font_size=13)

    add_callout_box(s, Cm(1.0), Cm(11.5), Cm(15.5), Cm(3.0),
                      "단조증가 패턴 일관 ★\n"
                      "no_shap < generic_rag < shaponly < fusion\n"
                      "→ 양 데이터셋 모두 동일 순서",
                      size=14, color=COLOR_HIGHLIGHT)
    add_callout_box(s, Cm(17.0), Cm(11.5), Cm(15.5), Cm(3.0),
                      "Value Match 0.59 → 0.90 (Home)\n"
                      "→ Fusion이 컨텍스트 값을 가장 정확히 인용\n"
                      "fact-grounded faithfulness 정량 입증",
                      size=14, color=COLOR_PRIMARY)
    add_footer(s, SW, SH, 17)

    # ═══════════════════════════════════════════════════
    # Slide 18 — G-Eval Cross-judge
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "16. G-Eval Cross-judge",
                       "Claude × Gemini 양방향 평가 — Completeness 일관, "
                       "Factual은 judge bias 입증")

    cj_path = FIG_DIR / "32_cross_judge_geval.png"
    if cj_path.exists():
        add_picture_centered(s, cj_path, top_cm=3.3, height_cm=9.5,
                                slide_width=SW)

    add_callout_box(s, Cm(1.0), Cm(13.2), Cm(15.5), Cm(3.8),
                      "Completeness — 양 judge 일관 양수\n"
                      "Fusion +0.67 (Claude judge) ~ +1.10 (Gemini judge)",
                      size=13, color=COLOR_HIGHLIGHT)
    add_callout_box(s, Cm(17.0), Cm(13.2), Cm(15.5), Cm(3.8),
                      "Factual Gemini target 차이 0.60 ★\n"
                      "Claude judge -0.13 vs Gemini judge +0.47\n"
                      "→ Self-preference bias 직접 입증",
                      size=13, color=COLOR_ACCENT)
    add_footer(s, SW, SH, 18)

    # ═══════════════════════════════════════════════════
    # Slide 19 — ★ Persona Trade-off
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "17. ★ Persona Trade-off — Honest Reporting",
                       "사실성 = fusion, 친근함 = generic_rag")

    pp_path = FIG_DIR / "36_human_proxy_personas.png"
    if pp_path.exists():
        add_picture_centered(s, pp_path, top_cm=3.2, height_cm=8.5,
                                slide_width=SW)

    add_textbox(s, Cm(1.0), Cm(12.0), Cm(31), Cm(0.8),
                  "Persona × Mode 1위 정리 (5점 만점)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    pp_rows = [
        ["평가 차원", "1 위", "2 위", "★ 비고"],
        ["사실성 (NLI)", "fusion 0.625", "shaponly 0.461", "본 연구 강점"],
        ["충실성 (G-Eval)", "fusion 4.82", "generic_rag 4.48", "Home 기준"],
        ["사람 친근성 (Persona trust)",
            "generic_rag 4.91", "fusion 4.31", "Fusion 2 위"],
        ["Customer clarity ⚠",
            "generic_rag 4.93", "fusion 2.67", "Fusion 약점 — honest"],
    ]
    add_table(s, Cm(1.5), Cm(13.0),
                [Cm(9), Cm(7), Cm(7), Cm(8.5)],
                [Cm(0.9)] * len(pp_rows), pp_rows, font_size=11)
    add_footer(s, SW, SH, 19)

    # ═══════════════════════════════════════════════════
    # Slide 20 — 공정성 Reweighing 4/4 통과
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "18. 공정성 Reweighing — 4 / 4 조합 통과",
                       "AUROC 손실 0.004 이내, 4/5 규칙 모두 통과")

    rw_path = FIG_DIR / "34_mitigation_bars.png"
    if rw_path.exists():
        add_picture_centered(s, rw_path, top_cm=3.3, height_cm=8.5,
                                slide_width=SW)

    rw_rows = [
        ["조합", "Before DI", "After DI", "AUROC Δ", "4/5 rule"],
        ["GENDER × Baseline", "0.62", "0.90", "−0.003", "✓ 통과"],
        ["GENDER × Aux", "0.64", "0.87", "+0.003", "✓ 통과"],
        ["AGE × Baseline", "0.56", "0.90", "−0.003", "✓ 통과"],
        ["AGE × Aux", "0.57", "0.83", "+0.003", "✓ 통과"],
    ]
    add_table(s, Cm(2.0), Cm(12.2),
                [Cm(9), Cm(5), Cm(5), Cm(6), Cm(5)],
                [Cm(0.9)] * len(rw_rows), rw_rows, font_size=11)
    add_footer(s, SW, SH, 20)

    # ═══════════════════════════════════════════════════
    # Slide 21 — 일반화 검증 (Home vs German)
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "19. 일반화 검증 — Home Credit vs UCI German",
                       "메커니즘 일관성 입증 — 약점 #5 해소")

    gen_path = FIG_DIR / "41_generalization.png"
    if gen_path.exists():
        add_picture_centered(s, gen_path, top_cm=3.2, height_cm=8.5,
                                slide_width=SW)

    gen_rows = [
        ["지표", "Home Credit", "UCI German Credit", "결과"],
        ["AUROC (XGBoost 5-fold)", "0.7587", "0.7714", "일관"],
        ["SHAP × Attention ρ", "0.117", "0.114", "거의 동일 ★"],
        ["NLI Entailment fusion", "0.625", "0.711", "German에서 더 강함 ★"],
        ["환각률", "0 %", "0 %", "Hard constraints 견고"],
        ["G-Eval Completeness 1 위", "fusion (4.82)",
            "generic_rag (4.58)", "데이터 복잡도 의존 ⚠"],
    ]
    add_table(s, Cm(1.5), Cm(12.0),
                [Cm(7.5), Cm(6), Cm(7), Cm(9.5)],
                [Cm(0.85)] * len(gen_rows), gen_rows, font_size=11)
    add_footer(s, SW, SH, 21)

    # ═══════════════════════════════════════════════════
    # Slide 22 — 정량 결과 종합 (한 표로 모든 핵심 수치)
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "20. 정량 결과 종합",
                       "본 연구가 학술적으로 입증한 모든 핵심 수치 — 한 눈에")

    add_textbox(s, Cm(1.0), Cm(3.0), Cm(31), Cm(0.7),
                  "9 개 영역 × 양 데이터셋 정량 결과 종합표",
                  size=13, bold=True, color=COLOR_PRIMARY)

    result_rows = [
        ["영역", "지표", "Home Credit", "UCI German", "결과 해석"],
        ["예측 성능", "XGBoost AUROC",
            "0.7587", "0.7714", "Mainstream 표준 (0.70~0.80)"],
        ["해석 일관성", "SHAP × Attn ρ",
            "0.117", "0.114", "거의 동일 — 메커니즘 일반화 ★"],
        ["환각률 (차원 B)", "LLM 변수명 환각",
            "0 / 60", "0 / 60", "Hard Constraints 견고"],
        ["사실성 (차원 C)", "NLI Entailment (fusion)",
            "0.625", "0.711", "4-mode 1 위 / 단조증가 ★"],
        ["값 정확성", "Value Match (fusion)",
            "0.90", "0.88", "4-mode 1 위 (0.59 → 0.90)"],
        ["충실성", "G-Eval Completeness 1 위",
            "fusion 4.82", "generic_rag 4.58", "데이터 복잡도 의존 ⚠"],
        ["사람 친화성", "Persona Trust 1 위",
            "generic_rag 4.91", "—", "fusion 2 위 (4.31) — Trade-off ★"],
        ["공정성", "DI ratio (Reweighing)",
            "0.62 → 0.90", "—", "4/4 통과, AUROC Δ ±0.003"],
        ["일반화", "양 데이터셋 일관성",
            "기준", "메커니즘 일관", "★ 본 연구 contribution"],
    ]
    add_table(s, Cm(1.0), Cm(3.9),
                [Cm(5.0), Cm(7.0), Cm(5.8), Cm(5.8), Cm(8.0)],
                [Cm(0.95)] * len(result_rows), result_rows, font_size=10)
    add_footer(s, SW, SH, 22)

    # ═══════════════════════════════════════════════════
    # Slide 23 — 본 연구의 5 가지 학술적 기여 + 정량 입증
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "21. 본 연구의 5 가지 학술적 기여",
                       "각 기여를 정량 결과로 입증 — 도입 RQ(슬라이드 5)에 대한 답")

    add_textbox(s, Cm(1.0), Cm(3.0), Cm(31), Cm(0.7),
                  "C1 ~ C5 기여 정리 + 정량 입증 근거 매핑",
                  size=13, bold=True, color=COLOR_PRIMARY)

    contrib_rows = [
        ["기여", "내용", "정량 입증 근거"],
        ["C1",
            "Agreement-aware Fusion Context 신규 설계\n"
            "(agreed · shap_only · attention_only 3 그룹 분해)",
            "SHAP × Attn ρ 0.117 / 0.114 (보완성)\n"
            "Agreement mean = 2.12 (n=100)"],
        ["C2",
            "4-tier 다중 평가 프레임워크 구축\n"
            "(Rules + NLI + Cross-Judge G-Eval + Persona)",
            "NLI 0.625, G-Eval +0.67 Completeness\n"
            "Persona 3 차원 × 5 점 척도"],
        ["C3",
            "4-mode 비교로 차별성 입증\n"
            "(환각 차단 ≠ fact-grounded 정확성)",
            "환각 0 % (모든 mode) vs NLI 단조증가\n"
            "Value Match 0.59 → 0.90 (no_shap → fusion)"],
        ["C4",
            "두 데이터셋 일반화 검증\n"
            "(Home Credit + UCI German Credit)",
            "양 데이터셋 SHAP × Attn ρ 거의 동일\n"
            "NLI 4-mode 단조증가 패턴 일관"],
        ["C5",
            "Reweighing 무손실 공정성 보정\n"
            "(4 / 4 조합 모두 4/5 규칙 통과)",
            "DI 0.62 → 0.90, 0.56 → 0.90\n"
            "AUROC 손실 ≤ 0.004"],
    ]
    add_table(s, Cm(1.0), Cm(3.9),
                [Cm(2.0), Cm(15.0), Cm(14.5)],
                [Cm(1.1), Cm(2.2), Cm(2.2), Cm(2.2), Cm(2.2), Cm(2.2)],
                contrib_rows, font_size=11)
    add_footer(s, SW, SH, 23)

    # ═══════════════════════════════════════════════════
    # Slide 24 — Honest Reporting — 한계 4가지
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "22. Honest Reporting — 연구의 한계",
                       "학술 논문의 정직성 표준에 따라 명시적 한계 보고")

    limits = [
        ("L1", "표본 크기 30 인스턴스",
            "4-mode × 2 LLM × 30 = 240 explanations — 통계적 의의는 trends 수준\n"
            "→ Future work: 표본 100+ 확장, 정식 통계 검정 (Wilcoxon paired)"),
        ("L2", "정식 IRB 인간평가 부재",
            "현 평가는 LLM-proxy persona (Claude judge) — 실제 사용자 데이터 X\n"
            "→ Future work: IRB 간소판 + 다수 평가자 + Cohen's κ"),
        ("L3", "Fusion의 Customer Clarity 약점",
            "Customer persona 기준 fusion 2.67 vs generic_rag 4.93 (큰 격차)\n"
            "→ Future work: Customer-friendly 표현 정제, hybrid fusion 모드"),
        ("L4", "Hard Constraint 의존성",
            "환각 차단은 hard constraints에 의존 — 제거 시 baseline no_shap 45.5 % 환각 확인\n"
            "→ Constraint 강도와 일반화 가능성 추가 검증 필요"),
    ]
    for i, (num, title, desc) in enumerate(limits):
        y = Cm(3.3 + i * 3.0)
        add_textbox(s, Cm(1.0), y, Cm(2.0), Cm(1.5),
                      num, size=22, bold=True, color=COLOR_ACCENT)
        add_textbox(s, Cm(3.0), y, Cm(29), Cm(1.0),
                      title, size=16, bold=True, color=COLOR_PRIMARY)
        add_textbox(s, Cm(3.0), y + Cm(0.9), Cm(29), Cm(2.0),
                      desc, size=11, color=COLOR_DARK)
    add_footer(s, SW, SH, 24)

    # ═══════════════════════════════════════════════════
    # Slide 25 — 응용 시나리오별 Mode 선택 가이드
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "23. 응용 시나리오별 Mode 선택 가이드",
                       "본 연구가 제시하는 실무 적용 권고")

    scenarios = [
        ("Audit / Regulation",
            "사실성 우선, 추적 가능성 필수",
            "★ fusion 권장",
            "Agreement 라벨 + Value Match 0.90\n→ 감사·규제 대응 적합",
            COLOR_HIGHLIGHT),
        ("Customer-facing",
            "이해 용이성·친근함 우선",
            "★ generic_rag 권장",
            "Customer clarity 4.93\n→ 일반 차주 대상 안내·고객센터",
            COLOR_PRIMARY),
        ("복잡 도메인 (Home Credit급)",
            "고차원·다변수·다양 관계",
            "★ fusion 권장",
            "G-Eval Completeness 1위\n→ 정보량이 많을수록 fusion 유리",
            RGBColor(0xDD, 0x88, 0x52)),
        ("단순 도메인 (German Credit급)",
            "저차원·소표본",
            "★ generic_rag 권장",
            "G-Eval Completeness 1위\n→ 도메인 지식 chunks가 보완",
            RGBColor(0x8C, 0x6B, 0xB1)),
    ]
    box_w = Cm(15.5); box_h = Cm(6); gap = Cm(0.5)
    for i, (scenario, criterion, recommend, why, color) in enumerate(scenarios):
        col = i % 2
        row = i // 2
        x = Cm(1.0) + col * (box_w + gap)
        y = Cm(3.5) + row * (box_h + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Cm(0.4); tf.margin_top = Cm(0.3)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = scenario
        _set_korean_font(r1, size=15, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = "기준: " + criterion
        _set_korean_font(r2, size=11, color=COLOR_WHITE)
        p3 = tf.add_paragraph()
        r3 = p3.add_run(); r3.text = "\n" + recommend
        _set_korean_font(r3, size=14, bold=True, color=COLOR_WHITE)
        p4 = tf.add_paragraph()
        r4 = p4.add_run(); r4.text = "\n" + why
        _set_korean_font(r4, size=10, color=COLOR_WHITE)
    add_footer(s, SW, SH, 25)

    # ═══════════════════════════════════════════════════
    # Slide 26 — 향후 연구 방향
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "24. 향후 연구 방향",
                       "본 연구를 확장할 수 있는 4 가지 방향")

    add_bullets(s, Cm(1.0), Cm(3.3), Cm(31), Cm(14), [
        {"text": "정식 IRB 인간평가 (L2 한계 해소)", "bold": True,
            "color": COLOR_ACCENT, "sub": [
            "IRB 간소판 + 5점 척도 + 다수 평가자(N=20~30)",
            "Cohen's κ 평가자 간 신뢰도",
            "사실성 vs 친근성 trade-off의 현실 검증 — 약점 #2 완전 해소"]},
        {"text": "추가 데이터셋 — 도메인 복잡도 stratification",
            "bold": True, "color": COLOR_PRIMARY, "sub": [
            "Australian Credit (UCI, 690×14, 복잡도 중간)",
            "Lending Club (Kaggle, ~1.3M×150)",
            "도메인 복잡도 × Mode 우월성 정량 매핑"]},
        {"text": "3-way Ablation — Attention-only 단독 검증",
            "bold": True, "color": COLOR_HIGHLIGHT, "sub": [
            "현재 비교: no_shap / generic_rag / shaponly / fusion",
            "추가 비교: attention_only → fusion 우월성이 "
            "attention 기여인지 isolate"]},
        {"text": "Customer-friendly Fusion 정제",
            "bold": True, "color": RGBColor(0x8C, 0x6B, 0xB1), "sub": [
            "Step 5-E 시도(revert): 친근함 ↑ 사실성 ↓ 정량 입증됨",
            "Future: Hybrid 표현 — agreement 라벨 자연어화, "
            "기술 용어 직관 변환"]},
    ], size=12)
    add_footer(s, SW, SH, 26)

    # ═══════════════════════════════════════════════════
    # Slide 27 — Q&A + 감사
    # ═══════════════════════════════════════════════════
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(0), Cm(0), Cm(1.5), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_textbox(s, Cm(2.5), Cm(4.5), Cm(28), Cm(2.0),
                  "감사합니다",
                  size=54, bold=True, color=COLOR_PRIMARY, align="center")
    add_textbox(s, Cm(2.5), Cm(8.0), Cm(28), Cm(1.5),
                  "Q & A",
                  size=32, color=COLOR_ACCENT, align="center")

    add_textbox(s, Cm(2.5), Cm(12.5), Cm(28), Cm(3.0),
                  "오픈소스 코드 저장소\n"
                  "https://github.com/Tim-Green0/tabnet-xai-rag-credit\n\n"
                  "A70067 오현택  |  지도교수 박운상\n"
                  "서강대학교 AI · SW대학원 데이터사이언스 · 인공지능 전공",
                  size=14, color=COLOR_DARK, align="center")

    return prs


def main():
    prs = make_slides()
    out = PAPER_DIR / "thesis_slides.pptx"
    prs.save(str(out))
    print(f"[OK] thesis 슬라이드 생성 완료: {out}")
    print(f"     파일 크기: {out.stat().st_size / 1024:.1f} KB")
    print(f"     슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
