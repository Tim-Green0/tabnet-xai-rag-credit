"""미팅용 발표 슬라이드 (.pptx) 생성기.

대상: 지도교수 중간 미팅, 15~20분 발표.
형식: 16:9, 14 슬라이드, 한국어, Malgun Gothic.

흐름:
   1. 표지
   2. 연구 배경 (왜 이 연구가 필요한가)
   3. 연구 목표 + 차별점
   4. 프레임워크 개요 (4단계 파이프라인)
   5. 데이터 + EDA 핵심 인사이트
   6. 모델 비교 (5-fold CV)
   7. SHAP × Attention 일관성 (RQ2)
   8. 공정성 진단 + Mitigation
   9. XAI-RAG Demo (idx 59291)
  10. 정량 평가 (RQ3 — 환각률 0)
  11. ★ Counterfactual Baseline 비교 (Claude 45% → 0%)
  12. 종합 + 핵심 메시지
  13. 한계 + 향후 계획
  14. Q&A

산출:  paper/midterm_slides.pptx

실행:  PYTHONIOENCODING=utf-8 .venv/Scripts/python.exe -m src.gen_slides
"""
from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Inches, Pt
from lxml import etree

from src.utils import RESULTS_DIR

PAPER_DIR = Path("paper")
PAPER_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR = Path("figures")

# 색상 팔레트
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)   # 진한 네이비 (제목, 강조)
COLOR_ACCENT = RGBColor(0xC4, 0x4E, 0x52)    # 빨강 (강조)
COLOR_HIGHLIGHT = RGBColor(0x55, 0xA8, 0x68)  # 초록 (긍정 결과)
COLOR_SUB = RGBColor(0x55, 0x55, 0x55)       # 회색 (부제)
COLOR_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFA)


# ─────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────
def _set_korean_font(run, size: int = 20, bold: bool = False,
                       color: RGBColor | None = None,
                       font_name: str = "Malgun Gothic"):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # East Asian font 설정 — 한글이 안 깨지게
    rPr = run._r.get_or_add_rPr()
    # 기존 ea 제거
    for ea_old in rPr.findall(qn("a:ea")):
        rPr.remove(ea_old)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    # latin도 같이
    for la_old in rPr.findall(qn("a:latin")):
        rPr.remove(la_old)
    la = etree.SubElement(rPr, qn("a:latin"))
    la.set("typeface", font_name)


def add_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height,
                  text: str, size: int = 20, bold: bool = False,
                  color: RGBColor | None = None,
                  align: str = "left", anchor: str = "top"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    _set_korean_font(run, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, left, top, width, height, bullets: list,
                  size: int = 16, color: RGBColor | None = None,
                  bullet_color: RGBColor | None = None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if isinstance(b, dict):
            text = b["text"]
            sub = b.get("sub", [])
            b_bold = b.get("bold", False)
            b_color = b.get("color", color)
        else:
            text = b
            sub = []
            b_bold = False
            b_color = color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet symbol
        run = p.add_run()
        run.text = "•  " + text
        _set_korean_font(run, size=size, bold=b_bold, color=b_color)
        # sub bullets (한 단계 들여쓰기)
        for s in sub:
            sp = tf.add_paragraph()
            sp.level = 1
            sp.alignment = PP_ALIGN.LEFT
            sr = sp.add_run()
            sr.text = "    –  " + s
            _set_korean_font(sr, size=size - 2, color=COLOR_SUB)
    return tb


def add_picture_centered(slide, path: Path, top_cm: float, height_cm: float,
                            slide_width: int, prs: Presentation):
    """슬라이드 가로 중앙에 그림 배치."""
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), Cm(0), Cm(top_cm),
                                       height=Cm(height_cm))
    # 가로 중앙
    pic.left = (slide_width - pic.width) // 2
    return pic


def add_header_bar(slide, slide_w: int, title: str,
                     subtitle: str | None = None):
    """슬라이드 상단 — 짙은 띠 + 제목."""
    # 좌측 컬러 띠
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(0), Cm(0),
                                       Cm(0.5), Cm(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # 제목
    add_textbox(slide, Cm(1.0), Cm(0.4), Cm(30), Cm(1.0),
                  title, size=24, bold=True, color=COLOR_PRIMARY)
    if subtitle:
        add_textbox(slide, Cm(1.0), Cm(1.4), Cm(30), Cm(0.8),
                      subtitle, size=12, color=COLOR_SUB)

    # 하단 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(1.0), Cm(2.4),
                                       Cm(31.5), Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()


def add_footer(slide, slide_w: int, slide_h: int, page_no: int,
                  total: int):
    add_textbox(slide, Cm(1.0), Cm(slide_h / 360000 - 0.9),
                  Cm(20), Cm(0.6),
                  "TabNet × SHAP × LLM XAI-RAG | A70067 오현택",
                  size=9, color=COLOR_SUB)
    add_textbox(slide, Cm(slide_w / 360000 - 3.0), Cm(slide_h / 360000 - 0.9),
                  Cm(2), Cm(0.6), f"{page_no} / {total}",
                  size=9, color=COLOR_SUB, align="right")


# ─────────────────────────────────────────────────────────────
# 슬라이드 생성
# ─────────────────────────────────────────────────────────────
def make_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    SW_CM = SW / 360000  # ≈ 33.87
    SH_CM = SH / 360000  # ≈ 19.05
    TOTAL = 14

    # ════════════════════════════════════════════════════════
    # 슬라이드 1 — 표지
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    # 배경 띠
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(0), Cm(0), SW, Cm(SH_CM))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    bg.line.fill.background()

    # 좌측 띠
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(0), Cm(0), Cm(1.5), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_textbox(s, Cm(2.5), Cm(4.5), Cm(28), Cm(2.5),
                  "TabNet과 거대언어모델(LLM) 기반\nXAI-RAG를 활용한 설명 가능한 신용 평가",
                  size=32, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Cm(2.5), Cm(8.5), Cm(28), Cm(1.0),
                  "석사 논문 중간 보고  |  Step 1 (8일 작업 완료)",
                  size=18, color=COLOR_SUB)
    add_textbox(s, Cm(2.5), Cm(13.5), Cm(28), Cm(2.5),
                  "전공: 데이터사이언스 · 인공지능\n"
                  "학번 A70067  |  오현택\n"
                  "지도교수: 박운상\n"
                  "보고일: 2026-05-10",
                  size=14, color=COLOR_DARK)

    # ════════════════════════════════════════════════════════
    # 슬라이드 2 — 연구 배경
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "1. 연구 배경",
                       "왜 신용 평가에 \"설명 가능한 자연어\"가 필요한가")

    add_bullets(s, Cm(1.0), Cm(3.2), Cm(30), Cm(15), [
        {"text": "신용 평가 자동화의 \"두 얼굴\"", "bold": True,
            "color": COLOR_PRIMARY,
            "sub": [
                "XGBoost·LightGBM·TabNet 같은 모델이 부도 예측 성능을 크게 끌어올림",
                "그러나 내부 의사결정이 직관적이지 않은 \"블랙박스\" 특성"]},
        {"text": "제도적·윤리적 요구",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "국내 금융소비자보호법 / 신용정보법 — 자동화된 신용평가 설명 의무",
                "EU GDPR — \"Right to Explanation\" 보장",
                "성별·연령 기반 잠재적 편향 → 차별 가능성"]},
        {"text": "기존 XAI(SHAP, LIME)의 한계",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "변수 중요도 그래프, 기여도 수치 — 전문가 외엔 해석 어려움",
                "\"모델은 설명 가능해졌으나 사람에게는 여전히 설명되지 않는다\""]},
        {"text": "LLM의 가능성과 위험",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "데이터-투-텍스트 변환에 강함",
                "그러나 환각(Hallucination) 문제 — 신용평가는 정확성·법적 책임이 중요"]},
    ], size=15)

    # ════════════════════════════════════════════════════════
    # 슬라이드 3 — 연구 목표 + 차별점
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "2. 연구 목표와 차별점", "본 연구가 푸는 문제")

    add_textbox(s, Cm(1.0), Cm(3.0), Cm(31), Cm(1.2),
                  "세 가지를 동시에 달성하는 통합 파이프라인을 제안",
                  size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.5), Cm(4.4), Cm(30), Cm(4), [
        "예측 성능: TabNet으로 트리 SOTA에 준하는 분류 성능",
        "투명한 근거: SHAP + TabNet 어텐션의 일관성 분석",
        "자연어 설명: LLM이 SHAP 사실만 받아쓰는 XAI-RAG로 환각 차단",
    ], size=16)

    add_textbox(s, Cm(1.0), Cm(10.0), Cm(31), Cm(1.0),
                  "기존 연구와 본 연구의 차별점", size=18, bold=True,
                  color=COLOR_PRIMARY)

    # 비교 표 (도형으로)
    comp_rows = [
        ["구분", "기존 연구", "본 연구"],
        ["LLM 활용", "분류기 또는 단순 변환기", "XAI-RAG로 환각 차단"],
        ["설명 평가", "정성적 사례", "Faithfulness · Hallucination · G-Eval 정량"],
        ["어텐션 vs SHAP", "암묵적 가정", "Spearman / Top-K 정량 비교"],
    ]
    table_top = Cm(11.2)
    col_widths = [Cm(7), Cm(11), Cm(13)]
    row_h = Cm(1.2)
    x_offset = Cm(1.0)
    for ri, row in enumerate(comp_rows):
        for ci, cell in enumerate(row):
            x = x_offset + sum(col_widths[:ci], Emu(0))
            y = table_top + ri * row_h
            shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                            col_widths[ci], row_h)
            if ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_PRIMARY
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = (RGBColor(0xF4, 0xF6, 0xFA)
                                                  if ri % 2 == 1 else
                                                  RGBColor(0xFF, 0xFF, 0xFF))
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.margin_top = Cm(0.1)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell
            _set_korean_font(run, size=12, bold=(ri == 0),
                                 color=(RGBColor(0xFF, 0xFF, 0xFF)
                                        if ri == 0 else COLOR_DARK))

    # ════════════════════════════════════════════════════════
    # 슬라이드 4 — 프레임워크
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "3. 프레임워크 개요",
                       "정형 데이터 → 예측 → 해석 → 컨텍스트 → 자연어 설명")

    # 4단계 파이프라인 — 박스 + 화살표
    stages = [
        ("① 예측", "TabNet / XGBoost\nP(default) 산출",
            RGBColor(0x4C, 0x72, 0xB0)),
        ("② 해석", "SHAP local +\nTabNet 어텐션",
            RGBColor(0xDD, 0x88, 0x52)),
        ("③ 컨텍스트", "JSON 사실 단위\n+ 민감변수 마스킹",
            RGBColor(0x55, 0xA8, 0x68)),
        ("④ 생성+평가", "LLM 자연어 설명\n+ Faithfulness/Halluc",
            RGBColor(0x8C, 0x6B, 0xB1)),
    ]
    box_w = Cm(7)
    box_h = Cm(4.5)
    gap = Cm(1.0)
    total_w = 4 * box_w + 3 * gap
    start_x = (SW - total_w) // 2
    box_y = Cm(4.5)
    for i, (label, desc, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_korean_font(run, size=20, bold=True,
                             color=RGBColor(0xFF, 0xFF, 0xFF))
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        _set_korean_font(r2, size=12, color=RGBColor(0xFF, 0xFF, 0xFF))

        # 화살표 (마지막 박스 제외)
        if i < 3:
            ax_x = x + box_w
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            ax_x + Cm(0.05),
                                            box_y + box_h // 2 - Cm(0.4),
                                            gap - Cm(0.1), Cm(0.8))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_SUB
            arrow.line.fill.background()

    add_textbox(s, Cm(1.0), Cm(11.0), Cm(31), Cm(1),
                  "핵심 — \"LLM은 SHAP 사실만 자연어로 변환하는 역할에 한정\"",
                  size=16, bold=True, color=COLOR_ACCENT, align="center")
    add_bullets(s, Cm(2.5), Cm(12.5), Cm(28), Cm(4), [
        "환각 원천 차단: 컨텍스트에 없는 변수·수치는 절대 생성 금지",
        "민감 변수 마스킹: 성별·연령 직접 노출 금지",
        "SHAP 부호 정확 반영: 양수는 부도↑, 음수는 부도↓",
    ], size=14)

    # ════════════════════════════════════════════════════════
    # 슬라이드 5 — 데이터 + EDA
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "4. 데이터와 EDA",
                       "Home Credit Default Risk (Kaggle) — 메인 테이블")

    add_bullets(s, Cm(1.0), Cm(3.0), Cm(15), Cm(15), [
        {"text": "307,511 행 × 122 컬럼", "bold": True, "color": COLOR_PRIMARY},
        {"text": "TARGET 분포 8.07% (불균형)",
            "sub": ["AUROC, AUPRC, KS 위주 평가"]},
        {"text": "EXT_SOURCE_2 / 3 가 단연 강한 신호",
            "sub": ["|ρ| 0.16~0.18, 다른 변수는 0.06 미만",
                     "→ 비선형/상호작용이 핵심"]},
        {"text": "결측 50%+ 컬럼 41개",
            "sub": ["*_MISSING_FLAG 변수 추가로 보존"]},
        {"text": "DAYS_EMPLOYED 365243 sentinel (18%)",
            "sub": ["NaN + EMPLOYED_FLAG 처리"]},
        {"text": "공정성 신호 이미 데이터에 존재",
            "color": COLOR_ACCENT, "bold": True,
            "sub": ["남성 부도율 1.45배, 25세 미만 vs 65세+ 3.4배"]},
    ], size=13)

    pic = s.shapes.add_picture(str(FIG_DIR / "03_protected_attrs.png"),
                                    Cm(17), Cm(3.5), Cm(15.5), Cm(11.5))
    add_textbox(s, Cm(17), Cm(15.2), Cm(15.5), Cm(0.6),
                  "그림. 성별·연령별 부도율 차이",
                  size=10, color=COLOR_SUB, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 6 — 모델 비교
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "5. 모델 비교",
                       "5-fold Stratified CV — test set, mean ± std")

    # 표
    cv_rows = [
        ["모델", "AUROC", "AUPRC", "KS", "특징"],
        ["XGBoost", "0.7587 ± 0.0008", "0.2445 ± 0.0011",
         "0.3846 ± 0.0015", "5/5 fold 1등"],
        ["LightGBM", "0.7544 ± 0.0009", "0.2402 ± 0.0018",
         "0.3788 ± 0.0028", "가장 빠름"],
        ["Logistic", "0.7544 ± 0.0001", "0.2343 ± 0.0006",
         "0.3804 ± 0.0010", "최고 안정성"],
        ["TabNet", "0.7518 ± 0.0017", "0.2331 ± 0.0023",
         "0.3749 ± 0.0056", "어텐션 해석성"],
    ]
    col_widths = [Cm(3.2), Cm(4.5), Cm(4.5), Cm(4.5), Cm(5.5)]
    row_h = Cm(1.0)
    table_top = Cm(3.2)
    x_offset = Cm(1.0)
    for ri, row in enumerate(cv_rows):
        for ci, cell in enumerate(row):
            x = x_offset + sum(col_widths[:ci], Emu(0))
            y = table_top + ri * row_h
            shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                            col_widths[ci], row_h)
            if ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_PRIMARY
                txt_color = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = (RGBColor(0xF4, 0xF6, 0xFA)
                                                  if ri % 2 == 1
                                                  else RGBColor(0xFF, 0xFF, 0xFF))
                txt_color = COLOR_DARK
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell
            _set_korean_font(run, size=11, bold=(ri == 0),
                                 color=txt_color)

    add_bullets(s, Cm(1.0), Cm(9.0), Cm(31), Cm(8), [
        {"text": "표준편차 0.001 안팎 — 단발 운이 아닌 안정적인 격차",
            "bold": True, "color": COLOR_HIGHLIGHT},
        {"text": "XGBoost가 미세하게 1등이지만 TabNet과 차이는 0.007 AUROC",
            "sub": ["성능보다 어텐션 해석성이 본 연구의 핵심 가치"]},
        {"text": "본 메시지: \"TabNet은 비교 가능한 성능 + 추가 해석성\"",
            "color": COLOR_PRIMARY, "bold": True},
    ], size=14)

    # ════════════════════════════════════════════════════════
    # 슬라이드 7 — SHAP × Attention
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "6. SHAP과 어텐션은 같은 것을 말하는가?",
                       "RQ2 — 두 해석 방법의 일관성 정량 분석")

    add_bullets(s, Cm(1.0), Cm(3.0), Cm(15), Cm(13), [
        {"text": "측정 방법", "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "TabNet 어텐션 importance",
                "SHAP global (KernelExplainer, 200 샘플)",
                "Spearman 상관 + Top-K 교집합/Jaccard"]},
        {"text": "결과 — 부분 일관 + 부분 상보", "bold": True,
            "color": COLOR_ACCENT,
            "sub": [
                "Spearman ρ (전체) = 0.117",
                "Spearman ρ (Top-50 합집합) = −0.195 (음수!)",
                "Top-20 교집합 = 9 / 20 (Jaccard 0.29)"]},
        {"text": "교집합 9개", "bold": True,
            "sub": [
                "EXT_SOURCE_2/3, DAYS_EMPLOYED, ORGANIZATION_TYPE,",
                "NAME_CONTRACT_TYPE_Revolving loans, CODE_GENDER_M 등"]},
        {"text": "결론 — 두 방법 함께 활용해야 입체적 해석",
            "bold": True, "color": COLOR_HIGHLIGHT},
    ], size=12)

    s.shapes.add_picture(str(FIG_DIR / "16_attention_vs_shap_scatter.png"),
                              Cm(17), Cm(3.5), height=Cm(11))
    add_textbox(s, Cm(17), Cm(15), Cm(15), Cm(0.6),
                  "그림. 어텐션 vs SHAP 산점도 (좌: 전체, 우: Top-20)",
                  size=10, color=COLOR_SUB, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 8 — 공정성
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "7. 공정성 진단 + Mitigation",
                       "DP / EO / EOdds / DI — 8개 케이스 모두 4/5 rule 위반")

    add_bullets(s, Cm(1.0), Cm(3.0), Cm(15), Cm(13), [
        {"text": "베이스라인 (4 모델 × {GENDER, AGE} = 8건)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "전부 DI < 0.8 (4/5 rule 위반)",
                "AGE DI ≈ 0.50, GENDER DI ≈ 0.62",
                "→ 알고리즘이 아닌 데이터 차원 편향"]},
        {"text": "Ablation (CODE_GENDER, DAYS_BIRTH 제거 후 재학습)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "GENDER ablation: DP −36~40% 감소 효과적",
                "AGE ablation: 거의 변화 없음 (proxy variable)",
                "AUROC 손실 1% 미만"]},
        {"text": "발견 — proxy variable 문제",
            "bold": True, "color": COLOR_ACCENT,
            "sub": [
                "연령 정보가 DAYS_EMPLOYED, OWN_CAR_AGE 등에",
                "간접 인코딩되어 단순 제거로 해결 안 됨",
                "→ Fairness-aware 학습이 future work로 자연스럽게 도출"]},
    ], size=12)

    s.shapes.add_picture(str(FIG_DIR / "19_fairness_mitigation.png"),
                              Cm(17), Cm(4.0), height=Cm(10))
    add_textbox(s, Cm(17), Cm(14.5), Cm(15), Cm(0.6),
                  "그림. baseline vs ablated AUROC + DP",
                  size=10, color=COLOR_SUB, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 9 — XAI-RAG Demo
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    walk = json.loads((RESULTS_DIR / "demo_walkthrough.json").read_text(encoding="utf-8"))
    sample = walk["sample"]
    pred = walk["prediction"]

    add_header_bar(s, SW, "8. XAI-RAG Demo (idx 59291)",
                       f"True Positive — 예측 P(default)={pred['default_proba']:.3f}, 결정 {pred['decision']}")

    add_textbox(s, Cm(1.0), Cm(3.2), Cm(31), Cm(0.8),
                  "한 명의 인스턴스가 4단계 파이프라인을 통과하는 흐름",
                  size=14, bold=True, color=COLOR_PRIMARY)

    # 좌: SHAP Top 5
    add_textbox(s, Cm(1.0), Cm(4.2), Cm(15), Cm(0.7),
                  "① SHAP Top 5 (거절 측 요인)",
                  size=12, bold=True, color=COLOR_DARK)
    shap_lines = []
    for d in walk["context"]["top_drivers_for_default"]:
        shap_lines.append(f"{d['rank']}. {d['feature']} = {d['value']}  (SHAP {d['shap']:+.3f})")
    add_textbox(s, Cm(1.2), Cm(5.0), Cm(14.5), Cm(5.5),
                  "\n".join(shap_lines),
                  size=10, color=COLOR_DARK)

    # 우: LLM 출력 (Claude만 — 더 자연스러움)
    add_textbox(s, Cm(17), Cm(4.2), Cm(15), Cm(0.7),
                  "② Claude Sonnet 4.5 자연어 설명 (발췌)",
                  size=12, bold=True, color=COLOR_DARK)
    claude_text = ""
    if "anthropic" in walk["llm_outputs"]:
        explanation = walk["llm_outputs"]["anthropic"].get("explanation", "")
        # 결정 요약 + 거절 사유 Top 3 부분만 추출
        lines = explanation.split("\n")
        keep = []
        section_count = 0
        for ln in lines:
            stripped = ln.strip()
            if not stripped:
                continue
            if "[" in stripped:
                section_count += 1
                if section_count > 2:
                    break
            keep.append(stripped[:90])
            if len(keep) > 12:
                break
        claude_text = "\n".join(keep)
    add_textbox(s, Cm(17.2), Cm(5.0), Cm(14.5), Cm(11),
                  claude_text or "(LLM 출력 없음)",
                  size=9, color=COLOR_DARK)

    # 하단 강조
    add_textbox(s, Cm(1.0), Cm(16.0), Cm(31), Cm(1),
                  "→ SHAP의 모든 변수·값·부호가 자연어 설명에 그대로 인용됨",
                  size=14, bold=True, color=COLOR_HIGHLIGHT, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 10 — 정량 평가 (RQ3)
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "9. 정량 평가 (RQ3)",
                       "Faithfulness · Hallucination · G-Eval — 두 LLM 모두에서 환각 0%")

    add_bullets(s, Cm(1.0), Cm(3.0), Cm(15), Cm(13), [
        {"text": "Hallucination Rate = 0.000",
            "bold": True, "color": COLOR_HIGHLIGHT,
            "sub": ["Gemini 2.5 Flash 0.000",
                     "Claude Sonnet 4.5 0.000"]},
        {"text": "G-Eval (Gemini self-judge, n=8)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "factual_accuracy 5.0 / 5",
                "sensitive_leak 5.0 / 5 (민감변수 마스킹 완벽)",
                "style 5.0 / 5",
                "completeness 3.4 / 5 (출력 schema 길이 제한 영향)"]},
        {"text": "Faithfulness 룰 기반",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "val_match Claude 0.901 vs Gemini 0.811",
                "sign_match Claude 0.867 vs Gemini 0.783"]},
        {"text": "효율성 (호출당)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": [
                "Claude: 8.4초 / 2,500 토큰",
                "Gemini: 12.7초 / 4,155 토큰 (reasoning ~60%)"]},
    ], size=12)

    s.shapes.add_picture(str(FIG_DIR / "21_llm_comparison.png"),
                              Cm(17), Cm(3.5), height=Cm(11))
    add_textbox(s, Cm(17), Cm(15), Cm(15), Cm(0.6),
                  "그림. Gemini vs Claude — 룰 기반 / G-Eval / 효율성",
                  size=10, color=COLOR_SUB, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 11 — Counterfactual Baseline ★
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "10. ★ 결정타 — SHAP이 없으면 환각이 늘어나는가?",
                       "동일 11 인스턴스 × \"SHAP 컨텍스트 없는\" baseline 실험")

    # 큰 숫자 강조 박스
    box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Cm(1.5), Cm(3.5), Cm(14), Cm(5.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_HIGHLIGHT
    box1.line.fill.background()
    add_textbox(s, Cm(2), Cm(4.0), Cm(13), Cm(1),
                  "XAI-RAG (SHAP context)", size=14, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(s, Cm(2), Cm(5.3), Cm(13), Cm(2.5),
                  "0.000", size=72, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
    add_textbox(s, Cm(2), Cm(8.0), Cm(13), Cm(0.8),
                  "Claude Sonnet 4.5",
                  size=12, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")

    box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Cm(17.5), Cm(3.5), Cm(14), Cm(5.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = COLOR_ACCENT
    box2.line.fill.background()
    add_textbox(s, Cm(18), Cm(4.0), Cm(13), Cm(1),
                  "Baseline (no SHAP)", size=14, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(s, Cm(18), Cm(5.3), Cm(13), Cm(2.5),
                  "45.5%", size=72, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
    add_textbox(s, Cm(18), Cm(8.0), Cm(13), Cm(0.8),
                  "Claude Sonnet 4.5",
                  size=12, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")

    add_textbox(s, Cm(1.0), Cm(10.5), Cm(31), Cm(1),
                  "Claude baseline 환각 사례",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(1.5), Cm(11.5), Cm(30), Cm(5), [
        "DTI · LTV · DSR — Home Credit에 없는 일반 금융 비율 약어",
        "\"햇살론, 미소금융\" — 데이터에 없는 한국 정부 지원 상품 추천",
        "\"☎ 1588-XXXX\" — 가짜 고객센터 번호 생성",
        "DEF_30_CNT — 변수명 잘림 (실제 DEF_30_CNT_SOCIAL_CIRCLE)",
    ], size=14)
    add_textbox(s, Cm(1.0), Cm(17.0), Cm(31), Cm(1),
                  "→ SHAP 컨텍스트가 LLM의 자유 추론을 직접 차단함을 입증",
                  size=14, bold=True, color=COLOR_ACCENT, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 12 (NEW) — Step 2-A 평가 신뢰성 강화: 100명 결과
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "11. ★ Step 2-A — 100명 표본에서도 환각 0%",
                       "Step 1의 약점(표본 10명) 해소 + 통계적 안정성 입증")

    # 좌측: 100명 결과
    add_textbox(s, Cm(1.0), Cm(3.2), Cm(15), Cm(1),
                  "표본 10배 확장 후 결과", size=18, bold=True,
                  color=COLOR_PRIMARY)
    box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Cm(1.0), Cm(4.3), Cm(15), Cm(7))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_HIGHLIGHT
    box1.line.fill.background()
    add_textbox(s, Cm(1.5), Cm(5.0), Cm(14), Cm(1),
                  "Hallucination Rate",
                  size=16, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
    add_textbox(s, Cm(1.5), Cm(6.2), Cm(14), Cm(2.5),
                  "0.000 / 100",
                  size=56, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                  align="center")
    add_textbox(s, Cm(1.5), Cm(9.0), Cm(14), Cm(1),
                  "Gemini 100건 + Claude 100건",
                  size=14, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")
    add_textbox(s, Cm(1.5), Cm(10.0), Cm(14), Cm(1),
                  "Step 1 (10건) 결과가 견고함을 통계로 입증",
                  size=12, color=RGBColor(0xFF, 0xFF, 0xFF), align="center")

    # 우측: Cross-LLM G-Eval 결과
    add_textbox(s, Cm(17.0), Cm(3.2), Cm(15), Cm(1),
                  "Cross-LLM G-Eval (self-bias 우회)",
                  size=18, bold=True, color=COLOR_PRIMARY)
    cross_rows = [
        ["Judge → Target", "Factual", "Sensitive", "Style"],
        ["Claude → Gemini", "4.87", "5.00", "4.97"],
        ["Gemini → Claude", "4.60", "5.00", "5.00"],
    ]
    col_widths_c = [Cm(4.5), Cm(3.5), Cm(3.5), Cm(3.5)]
    row_h_c = Cm(1.0)
    table_top_c = Cm(4.5)
    x_offset_c = Cm(17.0)
    for ri, row in enumerate(cross_rows):
        for ci, cell in enumerate(row):
            x = x_offset_c + sum(col_widths_c[:ci], Emu(0))
            y = table_top_c + ri * row_h_c
            shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                            col_widths_c[ci], row_h_c)
            if ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_PRIMARY
                txt_color = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = (RGBColor(0xF4, 0xF6, 0xFA)
                                                  if ri % 2 == 1
                                                  else RGBColor(0xFF, 0xFF, 0xFF))
                txt_color = COLOR_DARK
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = tf.paragraphs[0].add_run()
            run.text = cell
            _set_korean_font(run, size=11, bold=(ri == 0), color=txt_color)
    add_bullets(s, Cm(17.0), Cm(8.5), Cm(15), Cm(8), [
        "양 LLM 모두 factual ≥ 4.6/5",
        "sensitive_leak 5.0/5 만점 (양방향)",
        "Gemini self-bias가 자기 비판 방향",
    ], size=12)

    # ════════════════════════════════════════════════════════
    # 슬라이드 13 (NEW) — Counterfactual + Robustness
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "12. ★ Step 2-A — Counterfactual + Robustness",
                       "LLM이 컨텍스트에 의존하는가? 프롬프트 변형에 강건한가?")

    # 좌: Counterfactual
    add_textbox(s, Cm(1.0), Cm(3.0), Cm(15), Cm(1),
                  "Counterfactual Test (top driver 1개 제거)",
                  size=16, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Cm(1.0), Cm(4.0), Cm(15), Cm(1),
                  "원본 vs 변경 컨텍스트 출력 비교 (n=30)",
                  size=11, color=COLOR_SUB)
    cf_rows = [
        ["LLM", "Cosine sim", "ROUGE-L"],
        ["Claude Sonnet 4.5", "0.909 ± 0.069", "0.747"],
        ["Gemini 2.5 Flash", "0.920 ± 0.040", "0.750"],
    ]
    col_widths_cf = [Cm(5.5), Cm(5.0), Cm(4.5)]
    row_h_cf = Cm(1.0)
    table_top_cf = Cm(5.0)
    x_offset_cf = Cm(1.0)
    for ri, row in enumerate(cf_rows):
        for ci, cell in enumerate(row):
            x = x_offset_cf + sum(col_widths_cf[:ci], Emu(0))
            y = table_top_cf + ri * row_h_cf
            shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                            col_widths_cf[ci], row_h_cf)
            if ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_PRIMARY
                txt_color = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFA)
                txt_color = COLOR_DARK
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = tf.paragraphs[0].add_run()
            run.text = cell
            _set_korean_font(run, size=11, bold=(ri == 0), color=txt_color)
    add_bullets(s, Cm(1.0), Cm(8.5), Cm(15), Cm(8), [
        "cosine 0.91 → 의미는 90% 유지",
        "ROUGE-L 0.75 → 어휘는 25% 변경",
        "→ 부분 perturbation에 부분 반응 + 일관성 유지",
    ], size=12)

    # 우: Robustness
    add_textbox(s, Cm(17.0), Cm(3.0), Cm(15), Cm(1),
                  "Robustness (3 프롬프트 변형)",
                  size=16, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Cm(17.0), Cm(4.0), Cm(15), Cm(1),
                  "role/example/driver_shuffle (n=20)",
                  size=11, color=COLOR_SUB)
    rb_rows = [
        ["Variant", "Claude cosine", "Gemini cosine"],
        ["role_swap", "0.923", "0.951"],
        ["example_swap", "0.914", "0.908"],
        ["driver_shuffle", "0.924", "0.942"],
    ]
    col_widths_rb = [Cm(5.5), Cm(5.0), Cm(4.5)]
    row_h_rb = Cm(1.0)
    table_top_rb = Cm(5.0)
    x_offset_rb = Cm(17.0)
    for ri, row in enumerate(rb_rows):
        for ci, cell in enumerate(row):
            x = x_offset_rb + sum(col_widths_rb[:ci], Emu(0))
            y = table_top_rb + ri * row_h_rb
            shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                            col_widths_rb[ci], row_h_rb)
            if ri == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLOR_PRIMARY
                txt_color = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFA)
                txt_color = COLOR_DARK
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf = shape.text_frame
            tf.margin_left = Cm(0.2)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = tf.paragraphs[0].add_run()
            run.text = cell
            _set_korean_font(run, size=11, bold=(ri == 0), color=txt_color)
    add_bullets(s, Cm(17.0), Cm(9.5), Cm(15), Cm(8), [
        "양 LLM 모두 cosine ≥ 0.90 (목표 0.85+)",
        "Gemini가 약간 더 안정적",
        "→ 운영 환경에서도 출력 일관성 기대 가능",
    ], size=12)

    add_textbox(s, Cm(1.0), Cm(16.5), Cm(31), Cm(1),
                  "Step 1의 핵심 메시지(환각 0%, baseline 45.5%)가 통계적으로 견고함을 입증",
                  size=14, bold=True, color=COLOR_ACCENT, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 14 (NEW) — Step 3-B 보너스: 보조 테이블 활용
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "13. Step 3-B (보너스) — 보조 테이블 활용",
                       "발표 직전 D-5에 추가 진행한 성능 확장 (AUROC +2.22%)")

    # 좌측: 5-fold CV 비교 figure
    fig_path = FIG_DIR / "27_cv_aux_comparison.png"
    if fig_path.exists():
        pic = s.shapes.add_picture(str(fig_path),
                                       Cm(0.8), Cm(3.2),
                                       width=Cm(18.5))
        # 세로 비율 보존, 너무 크면 클립
        if pic.height > Cm(10):
            pic.height = Cm(10)
            pic.width = Cm(18.5)

    # 우측: 핵심 메시지
    add_textbox(s, Cm(20.0), Cm(3.2), Cm(13), Cm(1.2),
                  "AUROC 0.7587 → 0.7755 (+2.22%)",
                  size=18, bold=True, color=COLOR_ACCENT)
    add_bullets(s, Cm(20.3), Cm(4.5), Cm(13), Cm(7), [
        {"text": "추가한 보조 테이블 2개", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
                "bureau — 외부 신용기관 과거 대출 이력",
                "previous_application — Home Credit 자체 신청 이력"]},
        {"text": "5-fold CV 결과 (test set)", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
                "AUROC +0.0168 (baseline std의 21배)",
                "AUPRC +8.21%, KS +7.81% — 불균형 분류력 강화"]},
        {"text": "SHAP top 20 변화", "bold": True,
            "color": COLOR_PRIMARY, "sub": [
                "PREV_* feature 5개 신규 진입 (rank 12, 13, 14, 16, 20)",
                "최고 신호: 이전 거절 비율"]},
    ], size=11)

    add_textbox(s, Cm(1.0), Cm(15.5), Cm(32), Cm(2),
                  "★ 발견 — 외부 신용기관(bureau) feature는 top 20 진입 못함\n"
                  "  → main 테이블의 EXT_SOURCE_1/2/3에 외부 신용 정보가 응축돼 있을 가능성 (future work에서 ablation으로 검증)",
                  size=12, color=COLOR_HIGHLIGHT)

    # ════════════════════════════════════════════════════════
    # 슬라이드 15 (NEW) — Step 3-C-1: TabNet 어텐션 × SHAP 융합 컨텍스트
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "14. ★ Step 3-C-1 — TabNet 어텐션 × SHAP 융합 컨텍스트",
                       "TabNet이 본 메커니즘에 통합 — 환각 0% 유지 + Completeness 큰 향상")

    # 좌측: 비교 figure
    fig_path = FIG_DIR / "30_fusion_vs_shaponly.png"
    if fig_path.exists():
        pic = s.shapes.add_picture(str(fig_path),
                                       Cm(0.6), Cm(3.0),
                                       width=Cm(19))
        if pic.height > Cm(11):
            pic.height = Cm(11)
            pic.width = Cm(19)

    # 우측 상단: 메커니즘
    add_textbox(s, Cm(20.5), Cm(3.0), Cm(13), Cm(1.0),
                  "융합 메커니즘 (3 그룹)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(20.7), Cm(4.0), Cm(13), Cm(4), [
        "agreed_drivers: 두 모델 동의 강한 신호",
        "shap_only: SHAP만 본 보완 신호 (부호 보존)",
        "attention_only: TabNet만 본 sparse 신호 (부호 없음)",
        "그룹 라벨을 LLM에 명시 → 가중치 인지",
    ], size=11)

    # 우측 중간: agreement 통계
    add_textbox(s, Cm(20.5), Cm(8.0), Cm(13), Cm(1.0),
                  "Agreement 통계 (n=100)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(20.7), Cm(8.9), Cm(13), Cm(2.5), [
        "평균: agreed 2.12 / shap_only 6.98 / att_only 2.06",
        "3+ 동의 22%, 4+ 동의 0% — 부분 일관 + 상보",
    ], size=11)

    # 우측 하단: 평가 결과
    add_textbox(s, Cm(20.5), Cm(11.5), Cm(13), Cm(1.0),
                  "평가 결과 (n=30 each, judge=Claude)",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, Cm(20.7), Cm(12.4), Cm(13), Cm(4.5), [
        "Halluc 0/30 — 양 LLM × 양 mode 모두 ✅",
        "Completeness +0.67 (Anthropic) / +0.80 (Gemini)",
        "Factual 4.77~4.97, Sensitive 5.0/5.0 유지",
        "TabNet이 비교 모델 → 메커니즘 핵심으로 격상",
    ], size=11)

    # ════════════════════════════════════════════════════════
    # 슬라이드 16 (NEW) — Step 3-C-2: NLI 기반 평가 객관성 보강
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "15. ★ Step 3-C-2 — NLI로 평가 객관성 보강",
                       "룰 한계를 의미적 측정으로 입증 — Fusion이 NLI에서도 더 충실")

    # 좌측: NLI 비교 figure
    fig_path = FIG_DIR / "31_nli_vs_rules.png"
    if fig_path.exists():
        pic = s.shapes.add_picture(str(fig_path),
                                       Cm(0.6), Cm(3.0),
                                       width=Cm(19))
        if pic.height > Cm(11):
            pic.height = Cm(11)
            pic.width = Cm(19)

    # 우측 상단: 동기
    add_textbox(s, Cm(20.5), Cm(3.0), Cm(13), Cm(1.0),
                  "동기 — 룰 sign_match 하락 (Step 3-C-1)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Cm(20.7), Cm(4.0), Cm(13), Cm(3), [
        "Anthropic 0.87→0.65, Gemini 0.94→0.77",
        "원인 추정: 룰 키워드 셋의 한계 (다양한 표현 못 잡음)",
        "→ NLI로 의미적 함의를 직접 측정",
    ], size=11)

    # 우측 중간: 결과
    add_textbox(s, Cm(20.5), Cm(7.5), Cm(13), Cm(1.0),
                  "NLI 결과 (KLUE 다국어 NLI, n=30 each)",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, Cm(20.7), Cm(8.4), Cm(13), Cm(4.5), [
        "Anthropic: entailment +0.21 ★, contradiction -0.18 ★",
        "Gemini: entailment +0.12, contradiction -0.14",
        "양 LLM 일관 — fusion이 의미적으로 더 충실",
        "min_entailment(최악 문장) 양쪽 다 개선 또는 유지",
    ], size=11)

    # 우측 하단: 메시지
    add_textbox(s, Cm(20.5), Cm(13.5), Cm(13), Cm(3.5),
                  "결론 — 3-tier 평가 체계 완성\n"
                  "★ 룰의 sign_match 하락은 키워드 한계, 진짜 환각 아님 입증\n"
                  "★ Rules + G-Eval + NLI 다층 검증 → 평가 신뢰성 강화",
                  size=11, color=COLOR_HIGHLIGHT)

    # ════════════════════════════════════════════════════════
    # 슬라이드 17 — 종합 + 핵심 메시지 (기존 #15, 번호 갱신)
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "16. 종합 — 본 연구의 가치 4가지",
                       "Step 1 + 2-A + 3-B + 3-C-1 + 3-C-2 마무리")

    # 3개 컬럼 카드
    cards = [
        ("예측 성능",
          "XGBoost 0.7587 → 0.7755\n(Step 3-B aux 추가, +2.22%)\n\n"
          "TabNet 0.7518 — 어텐션\n해석성으로 융합 컨텍스트에 통합",
          RGBColor(0x4C, 0x72, 0xB0)),
        ("해석 융합",
          "ρ=0.117 부분 일관\n+ instance-level\n동의 22% (3+) / 0% (4+)\n\n"
          "agreement-aware 컨텍스트로\nLLM에 의미 라벨 전달",
          RGBColor(0x55, 0xA8, 0x68)),
        ("환각 차단",
          "Halluc 0/100 (Step 1/2-A)\n→ 0/30 (3-C-1 fusion)\n\n"
          "Baseline 45.5% vs\nXAI-RAG 0% (Step 1)\n"
          "NLI contradiction\n−0.14~−0.18 (3-C-2)",
          RGBColor(0xC4, 0x4E, 0x52)),
        ("완결성·충실성",
          "G-Eval Completeness\n+0.67~+0.80 (3-C-1)\n\n"
          "NLI Entailment\n+0.12~+0.21 (3-C-2)\n\n"
          "→ 다층 평가에서\n일관된 향상",
          RGBColor(0xDD, 0x85, 0x52)),
    ]
    card_w = Cm(7.0)
    card_h = Cm(8.5)
    gap_x = Cm(0.7)
    total_w = 4 * card_w + 3 * gap_x
    start_x = (SW - total_w) // 2
    card_y = Cm(3.5)
    for i, (title, body, color) in enumerate(cards):
        x = start_x + i * (card_w + gap_x)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.margin_left = Cm(0.6)
        tf.margin_right = Cm(0.6)
        tf.margin_top = Cm(0.6)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        _set_korean_font(run, size=22, bold=True,
                             color=RGBColor(0xFF, 0xFF, 0xFF))
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = "\n" + body
        _set_korean_font(r2, size=12, color=RGBColor(0xFF, 0xFF, 0xFF))

    add_textbox(s, Cm(1.0), Cm(13.0), Cm(31), Cm(1.5),
                  "한 줄 메시지",
                  size=16, bold=True, color=COLOR_PRIMARY, align="center")
    add_textbox(s, Cm(1.0), Cm(14.5), Cm(31), Cm(2),
                  "\"SHAP 컨텍스트를 LLM의 검색된 근거로 재정의하면,\n"
                  "LLM 종속성 없이 환각이 0이 되는 자연어 설명을 생성할 수 있다.\"",
                  size=18, bold=True, color=COLOR_DARK, align="center")

    # ════════════════════════════════════════════════════════
    # 슬라이드 15 — 한계 + 향후 계획
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    add_header_bar(s, SW, "17. 한계와 향후 계획",
                       "Step 3-B/3-C 일부 진행 / 잔여는 본 논문 확장")

    add_textbox(s, Cm(1.0), Cm(3.0), Cm(15), Cm(1),
                  "현재 한계 / 진행 중", size=16, bold=True, color=COLOR_ACCENT)
    add_bullets(s, Cm(1.5), Cm(4.0), Cm(15), Cm(13), [
        "Fusion 평가 표본 30명 — 100명+ 확장 검토",
        "Gemini judge cross-validation 미완 (503 과부하)",
        "인간평가 (Plausibility) 미수행 — IRB 절차 필요",
        "보조 테이블 2/6개 활용 (Step 3-B) — 4개 잔여",
        "Bureau의 SHAP 미진입 — EXT_SOURCE 응축 가설 검증 필요",
        "TabNet-only 컨텍스트 ablation 미수행 (3-way)",
        "Fairness-aware 학습 미수행",
    ], size=13)

    add_textbox(s, Cm(17), Cm(3.0), Cm(15), Cm(1),
                  "향후 계획", size=16, bold=True, color=COLOR_HIGHLIGHT)
    add_bullets(s, Cm(17.5), Cm(4.0), Cm(15), Cm(13), [
        "인간 평가 (Plausibility) — IRB + 5점 척도, Cohen's κ",
        "Gemini judge로 양방향 G-Eval cross-validation",
        "Fusion 표본 100명 확장 + counterfactual 결합",
        "3-way ablation: SHAP-only / Attn-only / Fusion",
        "잔여 보조 테이블 4개로 AUROC 0.78+",
        "Fairness-aware 학습 (Reweighing, Adversarial)",
        "FT-Transformer 비교 모델 추가",
    ], size=13)

    # ════════════════════════════════════════════════════════
    # 슬라이드 17 — Q&A
    # ════════════════════════════════════════════════════════
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(0), Cm(0), SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    add_textbox(s, Cm(1.0), Cm(7.0), Cm(31.5), Cm(3),
                  "감사합니다",
                  size=64, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                  align="center")
    add_textbox(s, Cm(1.0), Cm(11.0), Cm(31.5), Cm(1.5),
                  "Q & A",
                  size=36, color=RGBColor(0xCC, 0xDD, 0xEE),
                  align="center")
    add_textbox(s, Cm(1.0), Cm(13.5), Cm(31.5), Cm(1),
                  "GitHub: https://github.com/Tim-Green0/tabnet-xai-rag-credit",
                  size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
                  align="center")

    # 저장
    out = PAPER_DIR / "midterm_slides.pptx"
    prs.save(out)
    print(f"[OK] {out} 저장 ({len(prs.slides)} slides)")


if __name__ == "__main__":
    make_slides()
