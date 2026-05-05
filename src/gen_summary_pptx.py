"""연구 결과 소개 pptx 생성 (일회성, 날짜 prefix).

12 슬라이드 분량으로 본 연구의 흐름과 핵심 결과를 정리.
미팅용 발표 자료(midterm_slides.pptx)와 별개의 종합 소개 자료.

산출:
    paper/{YYYY-MM-DD}_연구결과정리.pptx

사용:
    PYTHONIOENCODING=utf-8 .venv/Scripts/python.exe -m src.gen_summary_pptx
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt
from lxml import etree

PAPER_DIR = Path("paper")
FIG_DIR = Path("figures")

# ─────────────────────────────────────────────────────────────
# 색상
# ─────────────────────────────────────────────────────────────
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)
COLOR_ACCENT = RGBColor(0xC4, 0x4E, 0x52)
COLOR_HIGHLIGHT = RGBColor(0x55, 0xA8, 0x68)
COLOR_SUB = RGBColor(0x55, 0x55, 0x55)
COLOR_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ─────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────
def _font(run, size=18, bold=False, color=None, name="Malgun Gothic"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ea_old in rPr.findall(qn("a:ea")):
        rPr.remove(ea_old)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)
    for la_old in rPr.findall(qn("a:latin")):
        rPr.remove(la_old)
    la = etree.SubElement(rPr, qn("a:latin"))
    la.set("typeface", name)


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, x_cm, y_cm, w_cm, h_cm, text,
                 size=14, bold=False, color=None, align=None):
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    if align == "center":
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    elif align == "right":
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        _font(run, size=size, bold=bold, color=color)
    return tb


def add_bullets(slide, x_cm, y_cm, w_cm, h_cm, items, size=14):
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, dict):
            text = "• " + item.get("text", "")
            run = p.add_run()
            run.text = text
            _font(run, size=size, bold=item.get("bold", False),
                   color=item.get("color"))
            for sub in item.get("sub", []):
                sp = tf.add_paragraph()
                sp.level = 1
                sr = sp.add_run()
                sr.text = "  – " + sub
                _font(sr, size=size - 1)
        else:
            run = p.add_run()
            run.text = "• " + item
            _font(run, size=size)
    return tb


def add_header(slide, slide_w, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Cm(0), Cm(0), slide_w, Cm(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    add_textbox(slide, 1.0, 0.4, slide_w / 360000 - 2.0, 1.4, title,
                  size=22, bold=True, color=COLOR_WHITE)
    if subtitle:
        add_textbox(slide, 1.0, 1.7, slide_w / 360000 - 2.0, 0.8, subtitle,
                      size=12, color=RGBColor(0xCC, 0xDD, 0xEE))


def add_image_bounded(slide, fig_path, x_cm, y_cm, max_w_cm, max_h_cm):
    """이미지를 비율 유지하며 max_w/max_h 안에 맞춤."""
    if not fig_path.exists():
        return None
    pic = slide.shapes.add_picture(str(fig_path), Cm(x_cm), Cm(y_cm),
                                       width=Cm(max_w_cm))
    if pic.height > Cm(max_h_cm):
        pic.height = Cm(max_h_cm)
        pic.width = Cm(max_w_cm)
    return pic


# ─────────────────────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────────────────────
def make_pptx() -> None:
    today = date.today().isoformat()
    out = PAPER_DIR / f"{today}_연구결과정리.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    SW_CM = SW / 360000

    # ────────── 슬라이드 1 — 표지 ──────────
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(1.5), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_textbox(s, 2.5, 4.0, 28, 2.5,
                  "TabNet × SHAP × LLM 기반\nXAI-RAG 신용 평가 — 연구 결과 정리",
                  size=30, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, 2.5, 8.5, 28, 1.0,
                  "환각 없는 자연어 설명을 위한 두 해석 신호 융합 RAG",
                  size=16, color=COLOR_SUB)
    add_textbox(s, 2.5, 13.0, 28, 3.0,
                  f"작성일: {today}\n"
                  "Step 1 (미팅 프로토타입) → Step 2-A (평가 신뢰성)\n"
                  "→ Step 3-B (성능 확장) → Step 3-C-1/2/2-f (TabNet 통합·NLI·Cross-Judge)\n"
                  "전공: 데이터사이언스 · 인공지능 석사  |  학번 A70067 · 오현택  |  지도교수: 박운상",
                  size=12, color=COLOR_DARK)

    # ────────── 슬라이드 2 — 한 줄 요약 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "1. 본 연구를 한 줄로",
                  "메커니즘 → 환각 차단의 직접 증거 → 다층 평가에서 일관 향상")
    add_textbox(s, 1.5, 4.0, 30, 3.0,
                  "\"SHAP과 TabNet 어텐션의 결합 결과를 LLM의 retrieved evidence로 재정의하면,\n"
                  "LLM 종속성 없이 환각이 0%인 자연어 설명 리포트를 생성할 수 있다.\"",
                  size=20, bold=True, color=COLOR_PRIMARY, align="center")
    add_bullets(s, 1.5, 9.0, 30, 8, [
        {"text": "정형 데이터 + 정확한 분류 (XGBoost AUROC 0.7755)",
            "bold": True, "color": COLOR_PRIMARY},
        {"text": "두 해석 신호의 부분 일관 + 부분 상보를 LLM에 명시 라벨로 전달",
            "bold": True, "color": COLOR_PRIMARY},
        {"text": "Halluc 0/100 + Counterfactual baseline 45.5% (no-SHAP) — 메커니즘 직접 증거",
            "bold": True, "color": COLOR_PRIMARY},
        {"text": "Rules + G-Eval(Cross-LLM) + NLI 3-tier × 양 LLM × 양 mode 다층 평가",
            "bold": True, "color": COLOR_PRIMARY},
    ], size=14)

    # ────────── 슬라이드 3 — 4단계 파이프라인 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "2. 4단계 파이프라인",
                  "정형 데이터 → 예측·해석 → 컨텍스트 → 자연어 설명 → 평가")

    boxes = [
        ("정형 데이터", "Home Credit\n307,511 × 122\n+ aux 756 features", 0x4C, 0x72, 0xB0),
        ("예측·해석", "XGBoost SHAP\n+ TabNet attention\n(local)", 0x55, 0xA8, 0x68),
        ("컨텍스트 빌더", "agreed / shap_only /\nattention_only\n(민감변수 마스킹)", 0xDD, 0x85, 0x52),
        ("LLM 자연어 설명", "Gemini 2.5 Flash\n+ Claude Sonnet 4.5\n(한국어)", 0x8C, 0x71, 0xB7),
        ("정량 평가", "Rules + G-Eval\n(Cross-LLM judge)\n+ NLI", 0xC4, 0x4E, 0x52),
    ]
    box_w = Cm(5.5)
    box_h = Cm(5.0)
    gap = Cm(0.7)
    total_w = 5 * box_w + 4 * gap
    start_x = (SW - total_w) // 2
    box_y = Cm(5.5)
    for i, (title, body, r, g, b) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, box_y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(r, g, b)
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Cm(0.4); tf.margin_top = Cm(0.4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        _font(run, size=14, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + body
        _font(r2, size=11, color=COLOR_WHITE)

        # 화살표 (마지막 박스 제외)
        if i < len(boxes) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + box_w, box_y + Cm(2.0),
                                            gap, Cm(1.0))
            arr.fill.solid(); arr.fill.fore_color.rgb = COLOR_PRIMARY
            arr.line.fill.background()

    # ────────── 슬라이드 4 — Step 1 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "3. Step 1 — 미팅용 작동 프로토타입",
                  "8일 작업으로 4단계 파이프라인 + Halluc 0% 검증")
    add_bullets(s, 1.5, 3.5, 31, 14, [
        {"text": "5-fold CV (test set) — XGBoost AUROC 0.7587 ± 0.0008 (1등)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["LightGBM 0.7549 / Logistic 0.7547 / TabNet 0.7543"]},
        {"text": "SHAP × TabNet 어텐션 일관성 (RQ2)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["Spearman ρ=0.117 (전체) / Top-50 ρ=−0.195 / Top-20 Jaccard=0.29",
                     "→ \"부분 일관 + 부분 상보\" 정량 입증 (Step 3-C-1 동기)"]},
        {"text": "공정성 진단 — 8/8 케이스 4/5 rule 위반 (DI < 0.8)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["GENDER ablation 효과적, AGE는 proxy로 간접 인코딩"]},
        {"text": "★ XAI-RAG 평가 (10 샘플) — Halluc 0/10 (양 LLM)",
            "bold": True, "color": COLOR_HIGHLIGHT},
        {"text": "★ Counterfactual Baseline — Step 1의 결정타",
            "bold": True, "color": COLOR_ACCENT,
            "sub": ["XAI-RAG (with SHAP context): Halluc 0%",
                     "no-SHAP baseline (raw 데이터만): Claude 45.5% — DTI/햇살론·미소금융 등 환각",
                     "→ SHAP 컨텍스트의 유무가 환각률에 결정적 차이를 만든다"]},
    ], size=12)

    # ────────── 슬라이드 5 — Step 2-A ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "4. Step 2-A — 평가 신뢰성 강화",
                  "100명 · Cross-LLM · Counterfactual · Robustness")
    add_bullets(s, 1.5, 3.5, 31, 14, [
        {"text": "★ 100명 표본 환각률 — 양 LLM 모두 0/100",
            "bold": True, "color": COLOR_HIGHLIGHT,
            "sub": ["Step 1의 10명 결과가 통계적으로 견고함을 입증"]},
        {"text": "Cross-LLM G-Eval 양방향 (n=30 each)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["Claude → Gemini: factual 4.87, completeness 4.0",
                     "Gemini → Claude: factual 4.6, completeness 3.33",
                     "양 LLM sensitive 5.0 만점 — 마스킹 정책 LLM 무관 견고"]},
        {"text": "Counterfactual 정량화 (top driver 1개 제거, n=30)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["Claude cosine 0.909, ROUGE-L 0.747",
                     "Gemini cosine 0.920, ROUGE-L 0.750",
                     "→ 부분 perturbation에 부분 반응 + 의미 일관 유지"]},
        {"text": "Robustness (3 변형 × n=20)",
            "bold": True, "color": COLOR_PRIMARY,
            "sub": ["cosine 0.908~0.951 — 프롬프트 미세 변형에 강건"]},
    ], size=12)

    # ────────── 슬라이드 6 — Step 3-B 성능 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "5. Step 3-B — 성능 확장 (보조 테이블)",
                  "bureau + previous_application 추가 → AUROC +2.22%")

    # 좌: figure
    add_image_bounded(s, FIG_DIR / "27_cv_aux_comparison.png",
                       0.6, 3.2, 19, 11)

    # 우: 메시지
    add_textbox(s, 20.5, 3.2, 13, 1.0,
                  "5-fold CV 결과 (test, mean ± std)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, 20.7, 4.2, 13, 7, [
        "AUROC 0.7587 → 0.7755 (+0.0168, +2.22%)",
        "AUPRC +8.21%, KS +7.81% — 분류력 자체 강화",
        "std 0.0011, 5 fold 모두 0.7743~0.7771 안정",
        "1161 features, 138s/fold — 운영 가능",
    ], size=11)

    add_textbox(s, 20.5, 11.5, 13, 1.0,
                  "★ SHAP top 20 변화 (의외)",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, 20.7, 12.5, 13, 4.5, [
        "신규 진입 5개 모두 PREV_* (자체 이력)",
        "최강: 이전 거절 비율 (PREV_NAME_CONTRACT_STATUS_Refused_mean)",
        "Bureau (외부 신용기관)는 진입 못함",
        "→ EXT_SOURCE에 응축됐을 가능성 (future work)",
    ], size=11)

    # ────────── 슬라이드 7 — Step 3-C-1 TabNet 통합 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "6. Step 3-C-1 — TabNet 어텐션 × SHAP 융합 컨텍스트",
                  "TabNet이 비교 모델 → 메커니즘 핵심으로 격상 (논문 제목 정당화)")

    add_image_bounded(s, FIG_DIR / "30_fusion_vs_shaponly.png",
                       0.6, 3.2, 19, 11)

    add_textbox(s, 20.5, 3.2, 13, 1.0,
                  "융합 메커니즘 (3 그룹)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, 20.7, 4.2, 13, 4, [
        "agreed: 두 모델 동의 강한 신호 (평균 2.12)",
        "shap_only: SHAP만 (부호 보존, 평균 6.98)",
        "attention_only: TabNet만 (sparse, 평균 2.06)",
        "n_agreed 분포 0~3, 4개+ 동의 0% — 부분 상보 입증",
    ], size=11)

    add_textbox(s, 20.5, 8.5, 13, 1.0,
                  "결과 (n=30 each, Claude judge)",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, 20.7, 9.5, 13, 7, [
        "Halluc 0/30 — 양 LLM × 양 mode ✅",
        "Completeness +0.67 (Anthropic) / +0.80 (Gemini)",
        "Factual 4.77~4.97 ≈ 유지",
        "Sensitive 5.0/5.0 만점 유지",
        "→ 환각 차단 + 완결성 향상 동시 달성",
    ], size=11)

    # ────────── 슬라이드 8 — Step 3-C-2 NLI ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "7. Step 3-C-2 — NLI로 평가 객관성 보강",
                  "룰 sign_match 하락이 키워드 한계임을 의미적 측정으로 입증")

    add_image_bounded(s, FIG_DIR / "31_nli_vs_rules.png",
                       0.6, 3.2, 19, 11)

    add_textbox(s, 20.5, 3.2, 13, 1.0,
                  "NLI 결과 (n=30 each)",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, 20.7, 4.2, 13, 5, [
        "Anthropic entailment +0.21 ★",
        "Anthropic contradiction −0.18 ★",
        "Gemini entailment +0.12, contradiction −0.14",
        "min_entailment(최악 문장) 양쪽 다 개선",
    ], size=11)

    add_textbox(s, 20.5, 9.0, 13, 1.0,
                  "★ 룰의 sign_match 하락 정체 입증",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, 20.7, 10.0, 13, 6.5, [
        "룰: 키워드(\"낮추는/긍정\") 매칭 → 한정 셋",
        "Fusion에서 LLM은 \"증가시키는, 위험\" 등 다양한 표현",
        "NLI는 의미적 함의 직접 측정 → 키워드 무관",
        "→ Fusion이 NLI에서도 더 충실 (양 LLM 일관)",
        "★ 3-tier: Rules + G-Eval + NLI 완성",
    ], size=11)

    # ────────── 슬라이드 9 — Step 3-C-2-f Cross-Judge ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "8. Step 3-C-2-f — Cross-Judge G-Eval (보너스)",
                  "Claude judge + Gemini judge 양방향으로 fusion 효과 검증")

    add_image_bounded(s, FIG_DIR / "32_cross_judge_geval.png",
                       0.6, 3.2, 19, 12)

    add_textbox(s, 20.5, 3.2, 13, 1.0,
                  "Δ = fusion - shaponly",
                  size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, 20.7, 4.2, 13, 6, [
        "Anthropic Compl.: Claude +0.67 / Gemini +0.90",
        "Gemini Compl.: Claude +0.80 / Gemini +1.10",
        "→ 양 judge 큰 향상 일관",
        "Sensitive 5.0/5.0 양 judge 만점",
    ], size=11)

    add_textbox(s, 20.5, 10.0, 13, 1.0,
                  "★ Cross-Judge의 가치 직접 증거",
                  size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, 20.7, 11.0, 13, 5.5, [
        "Gemini target Factual:",
        "  Claude judge −0.13 (사실상 동등)",
        "  Gemini judge +0.47 (명확한 향상)",
        "  차이 0.60 — 단일 judge였으면 잘못된 결론",
        "→ Cross-LLM judge가 본 연구에 필수",
    ], size=11)

    # ────────── 슬라이드 10 — 종합 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "9. 종합 — 본 연구의 5가지 가치", "")

    cards = [
        ("환각 차단",
          "Halluc 0/100\n(Step 1/2-A)\n+ 0/30 (3-C-1)\n\nBaseline 45.5%\nvs XAI-RAG 0%\n\nNLI contradiction\n−0.14~−0.18",
          RGBColor(0xC4, 0x4E, 0x52)),
        ("예측 성능",
          "XGBoost AUROC\n0.7587 → 0.7755\n(+2.22%, Step 3-B)\n\nAUPRC +8.21%\nKS +7.81%\n\n(불균형 분류력)",
          RGBColor(0x4C, 0x72, 0xB0)),
        ("해석 융합",
          "ρ=0.117 부분 일관\n+ instance-level\n부분 상보\n\nagreement-aware\n컨텍스트로\nLLM에 의미 라벨\n전달",
          RGBColor(0x55, 0xA8, 0x68)),
        ("완결성·충실성",
          "G-Eval Compl.\nClaude +0.67~+0.80\nGemini +0.90~+1.10\n\nNLI Entailment\n+0.12~+0.21\n\n다층 일관 향상",
          RGBColor(0xDD, 0x85, 0x52)),
        ("평가 신뢰성",
          "표본 10→100→30\n\nRules + G-Eval\n(Cross-LLM)\n+ NLI 3-tier\n\nRobustness\ncosine 0.91~0.95",
          RGBColor(0x8C, 0x71, 0xB7)),
    ]
    card_w = Cm(5.7)
    card_h = Cm(10)
    gap = Cm(0.4)
    total_w = 5 * card_w + 4 * gap
    start_x = (SW - total_w) // 2
    card_y = Cm(3.3)
    for i, (title, body, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, card_y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.margin_left = Cm(0.4); tf.margin_top = Cm(0.4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        _font(run, size=16, bold=True, color=COLOR_WHITE)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = "\n" + body
        _font(r2, size=10, color=COLOR_WHITE)

    add_textbox(s, 1.0, 14.5, SW_CM - 2.0, 2.0,
                  "★ 한 줄 — \"두 해석 신호의 상보성을 LLM에 명시 라벨로 제공해 사실성·민감변수·완결성을 손상시키지 않으면서 환각 0%를 유지하며, 다층 평가에서 일관된 향상을 보인다.\"",
                  size=12, bold=True, color=COLOR_DARK, align="center")

    # ────────── 슬라이드 11 — 한계와 향후 계획 ──────────
    s = add_blank_slide(prs)
    add_header(s, SW, "10. 한계와 향후 계획", "정직한 인식 + 미팅 후 우선순위")

    add_textbox(s, 1.0, 3.0, 15, 1.0,
                  "현재 한계",
                  size=15, bold=True, color=COLOR_ACCENT)
    add_bullets(s, 1.5, 4.0, 15, 12, [
        "Fusion 평가 표본 30 — 100+ 확장 검토",
        "인간평가 (Plausibility) 미수행 — 자동만으론 불충분",
        "보조 테이블 2/6개 — 4개 잔여",
        "Bureau의 SHAP 미진입 가설 ablation 미검증",
        "TabNet-only 컨텍스트 (3-way) 미수행",
        "Fairness-aware 학습 미수행",
        "Generic RAG baseline 미비교",
        "데이터 단일 (Home Credit) — 일반화 미검증",
        "한국어 native NLI 모델 추가 검증 필요",
    ], size=11)

    add_textbox(s, 17, 3.0, 15, 1.0,
                  "향후 계획 (우선순위)",
                  size=15, bold=True, color=COLOR_HIGHLIGHT)
    add_bullets(s, 17.5, 4.0, 15, 12, [
        "1순위: 인간평가 (IRB 간소판, 5점 척도, Cohen's κ)",
        "2순위: Fairness-aware 학습 (Reweighing 등)",
        "3순위: Generic RAG baseline 비교",
        "4순위: 잔여 보조 테이블 + Bureau ablation",
        "5순위: UCI German Credit 일반화",
        "6순위: 3-way ablation, TabNet/LightGBM 일반화",
        "장기: 한국어 도메인 특화 LLM (QLoRA)",
    ], size=11)

    # ────────── 슬라이드 12 — 마무리 ──────────
    s = add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    add_textbox(s, 1.0, 4.5, SW_CM - 2.0, 3.0,
                  "감사합니다",
                  size=64, bold=True, color=COLOR_WHITE, align="center")
    add_textbox(s, 1.0, 9.0, SW_CM - 2.0, 1.5,
                  f"연구 결과 정리 (Step 1 ~ Step 3-C-2-f)  |  {today}",
                  size=20, color=RGBColor(0xCC, 0xDD, 0xEE), align="center")
    add_textbox(s, 1.0, 12.0, SW_CM - 2.0, 1.5,
                  "GitHub: Tim-Green0/tabnet-xai-rag-credit",
                  size=14, color=RGBColor(0xAA, 0xCC, 0xEE), align="center")
    add_textbox(s, 1.0, 13.5, SW_CM - 2.0, 1.5,
                  "마일스톤: step1 → step2a → step3 → step4 (현재)",
                  size=12, color=RGBColor(0x99, 0xBB, 0xDD), align="center")

    PAPER_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"[OK] {out} 저장 ({len(prs.slides)} slides)")


if __name__ == "__main__":
    make_pptx()
